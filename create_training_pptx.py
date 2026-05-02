from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# カラーパレット
COLOR_MAIN = RGBColor(0x00, 0x5B, 0x99)       # 濃いブルー（タイトル背景）
COLOR_ACCENT = RGBColor(0x00, 0xA0, 0xC4)     # ライトブルー（アクセント）
COLOR_WARN = RGBColor(0xE8, 0x4C, 0x3B)       # レッド（警告・ポイント）
COLOR_BG = RGBColor(0xF5, 0xF8, 0xFF)         # 薄青白（背景）
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_SUBTEXT = RGBColor(0x44, 0x55, 0x66)
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)
COLOR_YELLOW = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=COLOR_DARK,
                align=PP_ALIGN.LEFT, wrap=True, font_name="メイリオ"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_slide_header(slide, title_text, subtitle_text="", slide_num=""):
    # 上部ヘッダーバー
    add_rect(slide, 0, 0, 13.33, 1.4, fill_color=COLOR_MAIN)
    # アクセントライン
    add_rect(slide, 0, 1.4, 13.33, 0.08, fill_color=COLOR_ACCENT)
    # 背景
    add_rect(slide, 0, 1.48, 13.33, 6.02, fill_color=COLOR_BG)

    # タイトル
    add_textbox(slide, title_text, 0.4, 0.15, 11.5, 0.9,
                font_size=28, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, subtitle_text, 0.4, 0.85, 10, 0.5,
                    font_size=14, color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.LEFT)
    if slide_num:
        add_textbox(slide, slide_num, 12.2, 0.2, 0.9, 0.5,
                    font_size=12, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.RIGHT)

def add_bullet_box(slide, items, left, top, width, height,
                   box_color=COLOR_WHITE, border_color=COLOR_ACCENT,
                   bullet="●", font_size=16, title=None, title_color=COLOR_MAIN):
    add_rect(slide, left, top, width, height, fill_color=box_color,
             line_color=border_color, line_width=Pt(1.5))
    y = top + 0.1
    if title:
        add_textbox(slide, title, left + 0.15, y, width - 0.3, 0.45,
                    font_size=16, bold=True, color=title_color)
        y += 0.45
    for item in items:
        add_textbox(slide, f"{bullet} {item}", left + 0.15, y, width - 0.3, 0.42,
                    font_size=font_size, color=COLOR_DARK)
        y += 0.42

def add_point_box(slide, text, left, top, width, height=0.6,
                  bg_color=COLOR_WARN, font_size=15, bold=True):
    add_rect(slide, left, top, width, height, fill_color=bg_color)
    add_textbox(slide, text, left + 0.1, top + 0.05, width - 0.2, height - 0.1,
                font_size=font_size, bold=bold, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

def add_image_placeholder(slide, left, top, width, height, label):
    add_rect(slide, left, top, width, height,
             fill_color=RGBColor(0xDD, 0xEA, 0xF5),
             line_color=COLOR_ACCENT, line_width=Pt(1.5))
    add_textbox(slide, "📷", left + width/2 - 0.3, top + height/2 - 0.35, 0.6, 0.6,
                font_size=28, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left + 0.1, top + height - 0.5, width - 0.2, 0.45,
                font_size=12, color=COLOR_SUBTEXT, align=PP_ALIGN.CENTER)


# ============================================================
# スライド 1 — タイトル
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=COLOR_MAIN)
add_rect(slide, 0, 5.5, 13.33, 2.0, fill_color=RGBColor(0x00, 0x45, 0x7A))
add_rect(slide, 0.5, 3.3, 12.33, 0.07, fill_color=COLOR_ACCENT)

add_textbox(slide, "新人職員研修", 1, 1.2, 11, 0.8,
            font_size=22, color=RGBColor(0xAA, 0xD4, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(slide, "移乗・移動介助の基本", 0.5, 2.0, 12.33, 1.2,
            font_size=46, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "〜安全・安楽・自立支援を実践する〜", 1, 3.5, 11, 0.7,
            font_size=20, color=RGBColor(0xCC, 0xE8, 0xFF), align=PP_ALIGN.CENTER)

add_textbox(slide, "対象：看護職員・介護職員\n所要時間：約30分", 1, 5.7, 11, 0.8,
            font_size=16, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(slide, "医療法人 研修・教育委員会", 1, 6.5, 11, 0.6,
            font_size=14, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


# ============================================================
# スライド 2 — 本日の研修目標
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "本日の研修目標", "この研修で身につけること", "2 / 18")

goals = [
    ("①", "移乗・移動介助の基本原則とボディメカニクスを理解する"),
    ("②", "車椅子・ベッド間の移乗手順を正しく実施できる"),
    ("③", "歩行介助・立ち上がり介助を安全に行える"),
    ("④", "転倒・転落リスクを予測し予防策を講じられる"),
    ("⑤", "患者の尊厳と残存能力を尊重した介助ができる"),
]

add_rect(slide, 0.5, 1.7, 12.33, 4.8, fill_color=COLOR_WHITE,
         line_color=COLOR_ACCENT, line_width=Pt(1.5))

for i, (num, text) in enumerate(goals):
    y = 1.85 + i * 0.88
    add_rect(slide, 0.6, y, 0.7, 0.65, fill_color=COLOR_MAIN)
    add_textbox(slide, num, 0.6, y + 0.05, 0.7, 0.55,
                font_size=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, text, 1.4, y + 0.05, 11.0, 0.6,
                font_size=17, color=COLOR_DARK)

add_point_box(slide, "患者の安全・安楽・自立支援が介助の3原則です",
              1.5, 6.6, 10.33, 0.6, bg_color=COLOR_ACCENT, font_size=16)


# ============================================================
# スライド 3 — 介助の3原則
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "介助の3原則", "すべての介助の土台となる考え方", "3 / 18")

principles = [
    ("安 全", COLOR_WARN,
     ["転倒・転落・皮膚損傷を防ぐ", "環境整備を必ず行う", "2人介助の基準を遵守する"]),
    ("安 楽", COLOR_ACCENT,
     ["患者の苦痛・不安を最小化", "声かけ・同意を必ず得る", "プライバシーを守る"]),
    ("自立支援", COLOR_GREEN,
     ["できることは自分で行ってもらう", "残存能力を活かす介助", "ADL維持・向上を目指す"]),
]

for i, (title, color, items) in enumerate(principles):
    x = 0.4 + i * 4.3
    add_rect(slide, x, 1.7, 4.0, 0.7, fill_color=color)
    add_textbox(slide, title, x, 1.75, 4.0, 0.6,
                font_size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 2.4, 4.0, 3.5, fill_color=COLOR_WHITE,
             line_color=color, line_width=Pt(2))
    for j, item in enumerate(items):
        add_textbox(slide, f"✓  {item}", x + 0.2, 2.55 + j * 0.75, 3.7, 0.65,
                    font_size=15, color=COLOR_DARK)

add_point_box(slide, "声かけ・同意・プライバシー保護は必ず実施！",
              1.5, 6.15, 10.33, 0.6, bg_color=COLOR_WARN)


# ============================================================
# スライド 4 — ボディメカニクスの基本
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "ボディメカニクスの基本", "介助者の腰痛予防と安全な介助のために", "4 / 18")

rules = [
    ("① 支持基底面を広げる", "足を肩幅に開き、安定した姿勢を確保する"),
    ("② 重心を低くする", "膝を曲げて腰を落とし、重心を下げる"),
    ("③ 重心を近づける", "患者との距離を縮め、密着して介助する"),
    ("④ 大きな筋群を使う", "腰だけでなく脚・体幹の大きな筋肉を使う"),
    ("⑤ 身体をひとつにまとめる", "患者の手足を体に引き付けてから動かす"),
    ("⑥ ねじり動作を避ける", "足を向けてから体を回旋させる"),
]

add_rect(slide, 0.4, 1.65, 8.5, 5.2, fill_color=COLOR_WHITE,
         line_color=COLOR_ACCENT, line_width=Pt(1.5))
for i, (title, desc) in enumerate(rules):
    y = 1.8 + i * 0.82
    add_rect(slide, 0.5, y, 3.2, 0.65, fill_color=COLOR_ACCENT)
    add_textbox(slide, title, 0.55, y + 0.05, 3.1, 0.55,
                font_size=14, bold=True, color=COLOR_WHITE)
    add_textbox(slide, desc, 3.85, y + 0.08, 4.9, 0.55,
                font_size=14, color=COLOR_DARK)

# 画像プレースホルダー
add_image_placeholder(slide, 9.1, 1.65, 3.8, 5.2,
                      "【画像①】ボディメカニクス6原則\nイラスト（姿勢比較図）")
add_point_box(slide, "腰痛は職業病ではありません。正しい姿勢で自分の体も守りましょう",
              0.5, 6.95, 12.33, 0.45, bg_color=COLOR_MAIN, font_size=14)


# ============================================================
# スライド 5 — 介助前の確認事項（リスクアセスメント）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "介助前の確認事項", "リスクアセスメントと環境整備", "5 / 18")

left_items = [
    "患者の理解・同意を得る",
    "バイタルサイン・体調確認",
    "点滴・ドレーン・カテーテル確認",
    "麻痺・疼痛・可動域の確認",
    "患者の体重・介助レベル確認",
]
right_items = [
    "ベッドの高さ調整（介助しやすい高さ）",
    "ブレーキ・ストッパーの確認",
    "床の濡れ・障害物の除去",
    "スリッパ→滑り止め靴に変更",
    "プライバシーカーテンを閉める",
]

add_bullet_box(slide, left_items, 0.4, 1.7, 6.1, 3.8,
               title="患者確認", title_color=COLOR_MAIN)
add_bullet_box(slide, right_items, 6.8, 1.7, 6.1, 3.8,
               title="環境整備", title_color=COLOR_GREEN)

add_point_box(slide, "⚠ 2人介助基準：全介助・体重60kg以上・不安定な患者は必ず2名で対応",
              0.4, 5.65, 12.53, 0.65, bg_color=COLOR_WARN, font_size=15)
add_textbox(slide, "介助の基本手順：確認 → 声かけ → 環境整備 → 介助 → 観察 → 記録",
            0.4, 6.38, 12.53, 0.5, font_size=15, bold=True,
            color=COLOR_MAIN, align=PP_ALIGN.CENTER)


# ============================================================
# スライド 6 — 立ち上がり介助
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "立ち上がり介助", "ベッド端座位から立位へ", "6 / 18")

steps = [
    ("STEP 1", "端座位を確保する", "ベッド端に座り、足底を床につける\n（足が浮いている場合はフットレストを使用）"),
    ("STEP 2", "前傾姿勢をとる", "体を前に傾け、重心を前方に移動\n介助者は腰部を支持する"),
    ("STEP 3", "立ち上がりを誘導", "「せーの」の声かけで一緒に立ち上がる\n患者の脇・腰を支えながら立位へ"),
    ("STEP 4", "立位の安定確認", "立位バランスを確認し、めまい・ふらつきがないか確認\n必要なら手すり・歩行器を使用"),
]

for i, (step, title, desc) in enumerate(steps):
    x = 0.35 + (i % 2) * 6.5
    y = 1.7 + (i // 2) * 2.5
    add_rect(slide, x, y, 6.2, 0.5, fill_color=COLOR_MAIN)
    add_textbox(slide, f"{step}  {title}", x + 0.1, y + 0.05, 6.0, 0.42,
                font_size=16, bold=True, color=COLOR_WHITE)
    add_rect(slide, x, y + 0.5, 6.2, 1.8, fill_color=COLOR_WHITE,
             line_color=COLOR_MAIN, line_width=Pt(1.2))
    add_textbox(slide, desc, x + 0.2, y + 0.6, 5.8, 1.6,
                font_size=14, color=COLOR_DARK)

add_point_box(slide, "ポイント：患者に合わせた速度で！急かすと転倒リスクが急増します",
              0.5, 6.85, 12.33, 0.5, bg_color=COLOR_ACCENT, font_size=14)


# ============================================================
# スライド 7 — ベッド→車椅子 移乗（手順）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "ベッド → 車椅子 移乗手順", "片麻痺患者（健側を活用した移乗）", "7 / 18")

steps7 = [
    "車椅子をベッドに対して約30°の角度で設置\n（健側に置く）",
    "車椅子のブレーキをかけ、フットレストを上げる",
    "患者を端座位にし、足底を床につける",
    "患者に前傾姿勢をとってもらい、健側の手で介助",
    "介助者は患側腰部〜臀部を支え、一緒に立ち上がる",
    "ゆっくり方向転換し、車椅子に向かって後ろに下がる",
    "臀部をシートに誘導しながらゆっくり着座させる",
    "フットレストを下げ、足を乗せて姿勢を整える",
]

add_rect(slide, 0.4, 1.65, 8.5, 5.25, fill_color=COLOR_WHITE,
         line_color=COLOR_ACCENT, line_width=Pt(1.5))
for i, step in enumerate(steps7):
    y = 1.78 + i * 0.62
    add_rect(slide, 0.5, y, 0.55, 0.48, fill_color=COLOR_ACCENT)
    add_textbox(slide, str(i + 1), 0.5, y + 0.03, 0.55, 0.42,
                font_size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, step, 1.15, y + 0.03, 7.6, 0.52,
                font_size=13, color=COLOR_DARK)

add_image_placeholder(slide, 9.1, 1.65, 3.8, 3.4,
                      "【画像②】ベッド→車椅子移乗\n手順イラスト（STEP1〜4）")
add_image_placeholder(slide, 9.1, 5.2, 3.8, 1.7,
                      "【画像③】車椅子の配置角度図\n（30°配置のイラスト）")


# ============================================================
# スライド 8 — 車椅子→ベッド 移乗
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "車椅子 → ベッド 移乗手順", "基本の流れと注意点", "8 / 18")

add_textbox(slide, "車椅子→ベッドは逆の手順で行います。ポイントを確認しましょう。",
            0.5, 1.7, 12.33, 0.5, font_size=16, color=COLOR_SUBTEXT)

points = [
    ("車椅子の配置", "ベッドに対して健側に30°で設置。ブレーキ・フットレスト確認を先に行う"),
    ("体重移動の誘導", "「前に体を倒して」と声かけ。前傾姿勢から立ち上がりを誘導"),
    ("方向転換", "介助者は患者の患側に位置。転倒防止のため体を密着させる"),
    ("着座誘導", "ベッドに手をつかせてから、ゆっくり膝を曲げて座らせる"),
    ("体位の整え", "座位バランス確認後、ベッド上での体位（臥位・座位）を整える"),
]

for i, (title, desc) in enumerate(points):
    y = 2.3 + i * 0.88
    add_rect(slide, 0.4, y, 2.8, 0.68, fill_color=COLOR_MAIN)
    add_textbox(slide, title, 0.45, y + 0.1, 2.7, 0.5,
                font_size=15, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.2, y, 9.7, 0.68, fill_color=COLOR_WHITE,
             line_color=COLOR_MAIN, line_width=Pt(1))
    add_textbox(slide, desc, 3.35, y + 0.1, 9.4, 0.55,
                font_size=14, color=COLOR_DARK)

add_point_box(slide, "両移乗共通：患者の「できること」を最大限活かす介助を心がけましょう",
              0.5, 6.8, 12.33, 0.55, bg_color=COLOR_GREEN, font_size=14)


# ============================================================
# スライド 9 — スライディングボード・介助用具
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "移乗補助用具の活用", "スライディングボード・移乗ベルト・リフト", "9 / 18")

tools = [
    ("スライディングボード",
     COLOR_ACCENT,
     ["座位移乗に使用", "車椅子⇔ベッド・シャワーチェア", "麻痺・筋力低下患者に有効", "使用時は必ず2名で確認"],
     "【画像④】スライディングボードの\n使用イラスト"),
    ("移乗ベルト（トランスファーベルト）",
     COLOR_GREEN,
     ["患者の腰に装着", "介助者がベルトを握って誘導", "急な体重移動に対応しやすい", "患者の転倒防止に有効"],
     "【画像⑤】移乗ベルト装着・\n使用方法イラスト"),
    ("介護リフト",
     COLOR_MAIN,
     ["全介助・重度障害患者に使用", "スリングシートを正しく装着", "事前にスリング選択が重要", "操作は必ず研修後に実施"],
     "【画像⑥】介護リフト操作\nイラスト"),
]

for i, (name, color, items, img_label) in enumerate(tools):
    x = 0.35 + i * 4.3
    add_rect(slide, x, 1.65, 4.0, 0.55, fill_color=color)
    add_textbox(slide, name, x + 0.1, 1.7, 3.8, 0.45,
                font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_image_placeholder(slide, x, 2.2, 4.0, 1.9, img_label)
    add_rect(slide, x, 4.1, 4.0, 2.5, fill_color=COLOR_WHITE,
             line_color=color, line_width=Pt(1.5))
    for j, item in enumerate(items):
        add_textbox(slide, f"・{item}", x + 0.15, 4.2 + j * 0.57, 3.7, 0.5,
                    font_size=13, color=COLOR_DARK)

add_point_box(slide, "用具の使用前は必ず破損・劣化チェックを行うこと",
              1.5, 6.7, 10.33, 0.55, bg_color=COLOR_WARN, font_size=14)


# ============================================================
# スライド 10 — 歩行介助の基本
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "歩行介助の基本", "安全な歩行介助の位置と手順", "10 / 18")

add_rect(slide, 0.4, 1.65, 12.53, 1.1, fill_color=COLOR_WHITE,
         line_color=COLOR_ACCENT, line_width=Pt(1))
add_textbox(slide,
            "介助者の位置：患側（麻痺のある側）の斜め後方  ／  片手：患者の腰または移乗ベルトを握る  ／  歩幅・速度は患者に合わせる",
            0.5, 1.75, 12.3, 0.85, font_size=15, bold=True, color=COLOR_MAIN, align=PP_ALIGN.CENTER)

walk_steps = [
    ("開始前", "靴・履物確認\nバイタル確認\n歩行補助具の用意"),
    ("出発時", "ベッドサイドで\nしばらく立位保持\nめまい確認"),
    ("歩行中", "声かけしながら\n一定ペースで\n障害物に注意"),
    ("方向転換", "健側に向かって\nゆっくり回転\n急な向き変え禁止"),
    ("着座時", "椅子/ベッドに\n確実に触れてから\nゆっくり着座"),
]

for i, (phase, desc) in enumerate(walk_steps):
    x = 0.35 + i * 2.52
    add_rect(slide, x, 2.9, 2.2, 0.55, fill_color=COLOR_MAIN)
    add_textbox(slide, phase, x, 2.95, 2.2, 0.45,
                font_size=15, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 3.45, 2.2, 1.8, fill_color=COLOR_WHITE,
             line_color=COLOR_MAIN, line_width=Pt(1))
    add_textbox(slide, desc, x + 0.1, 3.5, 2.0, 1.65, font_size=13, color=COLOR_DARK)

add_image_placeholder(slide, 0.5, 5.35, 5.5, 1.4,
                      "【画像⑦】歩行介助の位置・手の添え方イラスト")
add_image_placeholder(slide, 6.5, 5.35, 6.33, 1.4,
                      "【画像⑧】歩行補助具（T字杖・歩行器・ロフストランドクラッチ）の写真")


# ============================================================
# スライド 11 — 歩行補助具の種類
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "歩行補助具の種類と特徴", "患者の状態に合った補助具を選択する", "11 / 18")

tools11 = [
    ("T字杖（ステッキ）", COLOR_ACCENT,
     ["軽度の歩行障害に適応", "健側で使用する", "4点杖より安定性は低い"]),
    ("4点杖", COLOR_MAIN,
     ["接地面が広く安定", "片麻痺・バランス障害に有効", "室内使用が主"]),
    ("歩行器（固定型）", COLOR_GREEN,
     ["前方支持で安定", "両上肢に体重をかけられる", "段差・屋外には不向き"]),
    ("ロールウォーカー", COLOR_YELLOW,
     ["車輪付きで歩きやすい", "屋内外兼用可", "ブレーキ操作を確認"]),
]

for i, (name, color, items) in enumerate(tools11):
    x = 0.35 + i * 3.2
    add_rect(slide, x, 1.65, 3.0, 0.55, fill_color=color)
    add_textbox(slide, name, x + 0.05, 1.7, 2.9, 0.45,
                font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_image_placeholder(slide, x, 2.2, 3.0, 2.1, f"【写真】{name}")
    add_rect(slide, x, 4.3, 3.0, 2.1, fill_color=COLOR_WHITE,
             line_color=color, line_width=Pt(1.5))
    for j, item in enumerate(items):
        add_textbox(slide, f"・{item}", x + 0.1, 4.4 + j * 0.6, 2.8, 0.55,
                    font_size=13, color=COLOR_DARK)

add_point_box(slide, "補助具の選択は理学療法士・主治医と連携して決定すること",
              1.5, 6.5, 10.33, 0.6, bg_color=COLOR_MAIN, font_size=14)


# ============================================================
# スライド 12 — 全介助の移乗（2名介助）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "全介助の移乗（2名介助）", "意識障害・全身脱力・重度障害患者への対応", "12 / 18")

add_point_box(slide, "2名介助の適応：全介助患者 ／ 体重60kg以上 ／ 不安定な患者 ／ 初回移乗",
              0.4, 1.65, 12.53, 0.6, bg_color=COLOR_WARN, font_size=15)

roles = [
    ("リーダー役（頭側）", COLOR_MAIN, [
        "患者の上半身・頭頸部を支持",
        "声かけ・タイミングの指示を出す",
        "「いちにのさん」で全員同時に動く",
        "気道・カテーテル類の安全を確認",
    ]),
    ("サポート役（足側）", COLOR_ACCENT, [
        "患者の下肢・臀部を支持",
        "リーダーの指示に従う",
        "足先から骨盤まで確実に保持",
        "移動方向を確認してから動く",
    ]),
]

for i, (role, color, items) in enumerate(roles):
    x = 0.4 + i * 6.5
    add_rect(slide, x, 2.4, 6.1, 0.6, fill_color=color)
    add_textbox(slide, role, x + 0.1, 2.45, 5.9, 0.5,
                font_size=17, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 3.0, 6.1, 2.8, fill_color=COLOR_WHITE,
             line_color=color, line_width=Pt(1.5))
    for j, item in enumerate(items):
        add_textbox(slide, f"✓  {item}", x + 0.2, 3.1 + j * 0.65, 5.8, 0.58,
                    font_size=15, color=COLOR_DARK)

add_image_placeholder(slide, 0.4, 5.9, 12.53, 0.95,
                      "【画像⑨】2名介助でのベッド上移動・移乗イラスト（頭側・足側の持ち方）")


# ============================================================
# スライド 13 — 体位変換（臥位での移動）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "体位変換（臥位での移動）", "褥瘡予防と安楽な体位管理", "13 / 18")

add_rect(slide, 0.4, 1.65, 12.53, 0.6, fill_color=RGBColor(0xFF, 0xF3, 0xCD))
add_textbox(slide, "体位変換は原則2時間ごと。骨突出部（仙骨・踵・大転子等）への圧迫を避けること",
            0.5, 1.7, 12.3, 0.5, font_size=15, bold=True, color=RGBColor(0x85, 0x65, 0x04))

positions = [
    ("仰臥位\n（あお向け）", "最も基本の体位\n踵・仙骨の圧迫に注意"),
    ("側臥位\n（横向き）", "30°傾けが理想\nクッションで保持"),
    ("半座位\n（ファウラー位）", "床頭台15〜45°\n誤嚥予防に有効"),
    ("腹臥位\n（うつぶせ）", "呼吸改善に有効\n顔・腹部の圧迫注意"),
]

for i, (name, desc) in enumerate(positions):
    x = 0.4 + i * 3.2
    add_image_placeholder(slide, x, 2.35, 3.0, 1.9, f"【画像⑩-{i+1}】{name}の姿勢写真")
    add_rect(slide, x, 4.25, 3.0, 0.5, fill_color=COLOR_MAIN)
    add_textbox(slide, name, x, 4.28, 3.0, 0.45,
                font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 4.75, 3.0, 0.9, fill_color=COLOR_WHITE,
             line_color=COLOR_MAIN, line_width=Pt(1))
    add_textbox(slide, desc, x + 0.1, 4.8, 2.8, 0.8, font_size=13, color=COLOR_DARK)

add_point_box(slide, "体位変換のたびに皮膚状態・褥瘡チェックを必ず実施！",
              1.5, 5.75, 10.33, 0.6, bg_color=COLOR_WARN, font_size=14)

add_textbox(slide, "◎ポジショニングクッションを活用し骨突出部への圧迫を分散させましょう",
            0.5, 6.42, 12.33, 0.45, font_size=14, color=COLOR_MAIN, align=PP_ALIGN.CENTER)


# ============================================================
# スライド 14 — 転倒・転落予防
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "転倒・転落予防", "リスク評価と予防策の徹底", "14 / 18")

risks = [
    ("患者側リスク", COLOR_WARN, [
        "筋力低下・バランス障害",
        "認知症・見当識障害",
        "薬剤影響（睡眠薬・降圧薬等）",
        "起立性低血圧",
        "排泄に関する切迫感",
    ]),
    ("環境側リスク", COLOR_YELLOW, [
        "床の濡れ・障害物",
        "不適切な照明（暗すぎ）",
        "ベッド高さ不適切",
        "ナースコール未設置",
        "スリッパなど滑りやすい履物",
    ]),
    ("予防策", COLOR_GREEN, [
        "転倒リスクスコア評価（毎日）",
        "離床センサー・ベッドアラーム設置",
        "低床ベッドの使用",
        "ナースコール手の届く位置へ",
        "ラウンド強化（排泄パターン把握）",
    ]),
]

for i, (title, color, items) in enumerate(risks):
    x = 0.35 + i * 4.3
    add_rect(slide, x, 1.65, 4.0, 0.55, fill_color=color)
    add_textbox(slide, title, x, 1.7, 4.0, 0.45,
                font_size=17, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 2.2, 4.0, 3.8, fill_color=COLOR_WHITE,
             line_color=color, line_width=Pt(1.5))
    for j, item in enumerate(items):
        add_textbox(slide, f"▸ {item}", x + 0.15, 2.32 + j * 0.7, 3.75, 0.62,
                    font_size=14, color=COLOR_DARK)

add_point_box(slide, "転倒・転落発生時：患者の安全確保 → ナースコール → 状態確認 → 報告 → 記録",
              0.4, 6.1, 12.53, 0.65, bg_color=COLOR_WARN, font_size=14)


# ============================================================
# スライド 15 — よくあるミスと対策
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "よくあるミスとその対策", "新人職員が陥りやすいポイント", "15 / 18")

mistakes = [
    ("声かけを忘れる", "患者の心理的不安・筋緊張が増す", "必ず「〇〇しますね」と伝えてから開始"),
    ("ブレーキを忘れる", "移乗中に車椅子が動いて転落", "移乗前の確認を手順書通りに実施"),
    ("腰だけで持ち上げる", "介助者の腰痛・ぎっくり腰", "ボディメカニクスを意識。足・体幹を使う"),
    ("急ぎすぎる", "患者が不安・バランス崩す", "患者のペースに合わせ、ゆっくり確実に"),
    ("1人で無理をする", "患者・介助者ともに受傷リスク", "迷ったら必ず同僚に声をかける"),
    ("カテーテルを忘れる", "ドレーン抜去・皮膚損傷", "移乗前にルート類の長さ・固定を確認"),
]

add_rect(slide, 0.4, 1.65, 12.53, 5.1, fill_color=COLOR_WHITE,
         line_color=COLOR_ACCENT, line_width=Pt(1))

headers = ["よくあるミス", "起こりうる問題", "対策"]
widths = [3.0, 3.8, 5.3]
xs = [0.5, 3.6, 7.5]
for j, (h, w, x) in enumerate(zip(headers, widths, xs)):
    add_rect(slide, x, 1.7, w, 0.45, fill_color=COLOR_MAIN)
    add_textbox(slide, h, x, 1.72, w, 0.42,
                font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

for i, (m, prob, sol) in enumerate(mistakes):
    y = 2.2 + i * 0.72
    bg = RGBColor(0xF0, 0xF5, 0xFF) if i % 2 == 0 else COLOR_WHITE
    for x, w in zip(xs, widths):
        add_rect(slide, x, y, w, 0.65, fill_color=bg,
                 line_color=RGBColor(0xCC, 0xDD, 0xEE), line_width=Pt(0.5))
    texts = [f"✗ {m}", prob, f"✓ {sol}"]
    colors = [COLOR_WARN, RGBColor(0x80, 0x50, 0x00), COLOR_GREEN]
    for (txt, col, x, w) in zip(texts, colors, xs, widths):
        add_textbox(slide, txt, x + 0.1, y + 0.08, w - 0.2, 0.5,
                    font_size=12, color=col)

add_point_box(slide, "わからないことは必ず先輩に確認！報告・連絡・相談を徹底しましょう",
              1.5, 6.85, 10.33, 0.5, bg_color=COLOR_MAIN, font_size=14)


# ============================================================
# スライド 16 — 記録・報告のポイント
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "記録・報告のポイント", "介助後の観察と情報共有", "16 / 18")

record_items = [
    ("いつ", "移乗・移動の実施日時"),
    ("どこで", "場所（病室・廊下・リハビリ室など）"),
    ("何を", "実施した介助内容（車椅子移乗・歩行介助など）"),
    ("誰が", "実施者（1名・2名介助など）"),
    ("どのように", "使用した補助具・患者の協力度"),
    ("患者の状態", "バイタル変化・表情・疼痛・ふらつきの有無"),
    ("特記事項", "ヒヤリハット・異常所見は必ず記録・報告"),
]

add_rect(slide, 0.4, 1.65, 7.5, 5.3, fill_color=COLOR_WHITE,
         line_color=COLOR_ACCENT, line_width=Pt(1.5))
add_textbox(slide, "記録に含める内容（5W1H＋状態）",
            0.5, 1.7, 7.3, 0.5, font_size=16, bold=True, color=COLOR_MAIN)

for i, (key, val) in enumerate(record_items):
    y = 2.25 + i * 0.64
    add_rect(slide, 0.5, y, 1.5, 0.52, fill_color=COLOR_ACCENT)
    add_textbox(slide, key, 0.5, y + 0.06, 1.5, 0.42,
                font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, val, 2.1, y + 0.08, 5.7, 0.42, font_size=14, color=COLOR_DARK)

add_rect(slide, 8.1, 1.65, 4.8, 5.3, fill_color=COLOR_WHITE,
         line_color=COLOR_WARN, line_width=Pt(2))
add_rect(slide, 8.1, 1.65, 4.8, 0.55, fill_color=COLOR_WARN)
add_textbox(slide, "報告が必要な場合", 8.1, 1.7, 4.8, 0.45,
            font_size=15, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

report_cases = [
    "転倒・転落が発生した",
    "ヒヤリハットがあった",
    "患者が痛みを訴えた",
    "バイタルが基準値を逸脱した",
    "皮膚の発赤・損傷を発見した",
    "ドレーン・カテーテルが抜けた",
    "患者が介助を拒否した",
]
for i, case in enumerate(report_cases):
    add_textbox(slide, f"⚠ {case}", 8.2, 2.28 + i * 0.63, 4.6, 0.55,
                font_size=13, color=COLOR_WARN if i < 2 else COLOR_DARK)

add_point_box(slide, "記録は「事実」を具体的・客観的に。主観的な表現は避ける",
              0.4, 7.0, 12.53, 0.4, bg_color=COLOR_MAIN, font_size=13)


# ============================================================
# スライド 17 — 確認テスト
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "確認テスト", "本日の研修内容を振り返りましょう", "17 / 18")

questions = [
    ("Q1", "ボディメカニクスの目的は何ですか？\n①患者の安楽  ②介助者の腰痛予防  ③両方  ④どちらでもない"),
    ("Q2", "片麻痺患者への車椅子移乗で、車椅子を置く位置はどちら側ですか？\n①患側（麻痺側）  ②健側（麻痺のない側）  ③どちらでも同じ"),
    ("Q3", "2名介助が必要な条件を1つ答えてください。"),
    ("Q4", "移乗介助前に必ず確認すべき3つのことを答えてください。"),
    ("Q5", "転倒が発生した際の対応の優先順位を答えてください。"),
]

answers = ["③両方", "②健側", "全介助・体重60kg以上・不安定な患者など",
           "患者の同意・ブレーキ・ルート類など", "安全確保→ナース→状態確認→報告→記録"]

for i, (qnum, qtext) in enumerate(questions):
    y = 1.7 + i * 1.02
    add_rect(slide, 0.4, y, 1.0, 0.75, fill_color=COLOR_MAIN)
    add_textbox(slide, qnum, 0.4, y + 0.1, 1.0, 0.55,
                font_size=16, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.4, y, 11.53, 0.75, fill_color=COLOR_WHITE,
             line_color=COLOR_MAIN, line_width=Pt(1))
    add_textbox(slide, qtext, 1.55, y + 0.08, 11.2, 0.65, font_size=13, color=COLOR_DARK)

add_point_box(slide, "答え合わせは研修担当者と一緒に行いましょう！",
              2.0, 6.85, 9.33, 0.5, bg_color=COLOR_ACCENT, font_size=14)


# ============================================================
# スライド 18 — まとめ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide, "本日のまとめ", "移乗・移動介助のポイントを復習", "18 / 18")

summary = [
    ("1", "介助の3原則", "安全・安楽・自立支援を常に意識する"),
    ("2", "ボディメカニクス", "正しい姿勢で患者も介助者も守る"),
    ("3", "確認・声かけ", "介助前の確認と患者への声かけを徹底する"),
    ("4", "移乗手順の遵守", "車椅子の角度・ブレーキ・フットレストを確認"),
    ("5", "2名介助の判断", "一人で無理せず、迷ったら必ず声をかける"),
    ("6", "記録・報告", "気づいたことはすぐに記録・報告・連絡"),
]

for i, (num, title, desc) in enumerate(summary):
    x = 0.35 if i < 3 else 6.7
    y = 1.7 + (i % 3) * 1.48
    add_rect(slide, x, y, 0.7, 1.2, fill_color=COLOR_MAIN)
    add_textbox(slide, num, x, y + 0.3, 0.7, 0.6,
                font_size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.7, y, 5.65, 1.2, fill_color=COLOR_WHITE,
             line_color=COLOR_MAIN, line_width=Pt(1.5))
    add_textbox(slide, title, x + 0.85, y + 0.1, 5.4, 0.45,
                font_size=16, bold=True, color=COLOR_MAIN)
    add_textbox(slide, desc, x + 0.85, y + 0.55, 5.4, 0.55,
                font_size=14, color=COLOR_DARK)

add_rect(slide, 0.4, 6.6, 12.53, 0.75, fill_color=COLOR_MAIN)
add_textbox(slide, "患者の「安全・安楽・自立支援」のために、今日から実践しましょう！",
            0.5, 6.68, 12.3, 0.55,
            font_size=17, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


# 保存
output_path = r"C:\Users\yuuna\agens\移乗移動介助_新人研修.pptx"
prs.save(output_path)
print(f"保存完了: {output_path}")
print(f"スライド数: {len(prs.slides)}枚")
