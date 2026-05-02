from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ===== カラーパレット =====
# ベース
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # ほぼ黒（メインテキスト）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY   = RGBColor(0xF4, 0xF6, 0xF9)   # 薄グレー背景
C_MGRAY   = RGBColor(0xB0, 0xBE, 0xC5)   # 中グレー（サブテキスト）

# セクションカラー
C_BLUE    = RGBColor(0x1B, 0x4F, 0x8A)   # 個人情報 — ネイビーブルー
C_LBLUE   = RGBColor(0xDB, 0xEA, 0xF8)   # 個人情報 — 薄ブルー
C_TEAL    = RGBColor(0x0D, 0x73, 0x77)   # 法令遵守 — ティール
C_LTEAL   = RGBColor(0xD4, 0xEF, 0xED)   # 法令遵守 — 薄ティール
C_ORANGE  = RGBColor(0xE6, 0x7E, 0x22)   # ハラスメント — オレンジ
C_LORANGE = RGBColor(0xFD, 0xF0, 0xE0)   # ハラスメント — 薄オレンジ
C_RED     = RGBColor(0xC0, 0x39, 0x2B)   # カスハラ・違反 — レッド
C_LRED    = RGBColor(0xFA, 0xDB, 0xD7)   # カスハラ — 薄レッド
C_GOLD    = RGBColor(0xF3, 0x9C, 0x12)   # アクセント — ゴールド
C_GREEN   = RGBColor(0x1E, 0x8B, 0x4C)   # OK系 — グリーン

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ===== ユーティリティ =====

def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def tx(slide, text, x, y, w, h, size=14, bold=False,
       color=C_DARK, align=PP_ALIGN.LEFT, wrap=True, font='Meiryo'):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    try:
        run._r.rPr.set(qn('w:eastAsia'), font)
    except Exception:
        pass
    return txb, tf

def add_para(tf, text, size=12, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, sp=Pt(3), font='Meiryo'):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = sp
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return p

def slide_no(slide, num, total=13):
    tx(slide, f'{num} / {total}',
       SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.38),
       Inches(0.8), Inches(0.3),
       size=9, color=C_MGRAY, align=PP_ALIGN.RIGHT)

def header_bar(slide, accent_color, section_tag='', title=''):
    """スライド上部カラーバー + タイトル"""
    rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=accent_color)
    if title:
        tx(slide, title,
           Inches(0.5), Inches(0.18), Inches(9.5), Inches(0.62),
           size=22, bold=True, color=C_DARK)
    if section_tag:
        tx(slide, section_tag,
           SLIDE_W - Inches(3.5), Inches(0.18), Inches(3.3), Inches(0.35),
           size=8, color=accent_color, align=PP_ALIGN.RIGHT)
    # 下線
    rect(slide, Inches(0.5), Inches(0.82), Inches(12.3), Pt(1.5),
         fill=accent_color)

def footer_bar(slide, accent_color, msg=''):
    rect(slide, 0, SLIDE_H - Inches(0.42), SLIDE_W, Inches(0.42), fill=accent_color)
    if msg:
        tx(slide, msg,
           Inches(0.5), SLIDE_H - Inches(0.38), Inches(12.0), Inches(0.35),
           size=10, bold=True, color=C_WHITE)


# ============================================================
# スライド 1: 表紙
# ============================================================
s = prs.slides.add_slide(BLANK)

# 左ブルー帯
rect(s, 0, 0, Inches(5.5), SLIDE_H, fill=C_BLUE)
# 左帯アクセント（下端）
rect(s, 0, SLIDE_H - Inches(0.5), Inches(5.5), Inches(0.5), fill=C_GOLD)

tx(s, '2026年度　院内必修研修',
   Inches(0.4), Inches(0.9), Inches(4.8), Inches(0.4),
   size=11, color=C_MGRAY)

for i, line in enumerate(['個人情報保護', 'ハラスメント予防', '法令遵守']):
    tx(s, line,
       Inches(0.4), Inches(1.55 + i * 1.1), Inches(4.9), Inches(1.0),
       size=30, bold=True, color=C_WHITE)

rect(s, Inches(0.4), Inches(5.0), Inches(1.5), Pt(3), fill=C_GOLD)
tx(s, 'すべての病院スタッフへ',
   Inches(0.4), Inches(5.15), Inches(4.5), Inches(0.4),
   size=13, color=C_MGRAY)
tx(s, '研修時間：約30分　／　対象：全職員',
   Inches(0.4), Inches(5.65), Inches(4.5), Inches(0.35),
   size=10, color=RGBColor(0x80, 0x90, 0xA0))

# 右　テーマバッジ
badge_data = [
    ('📋  個人情報保護',  C_BLUE,   C_LBLUE,  '患者情報の保護と漏洩防止'),
    ('🛡  ハラスメント予防', C_ORANGE, C_LORANGE, '職場内・カスハラ重点解説'),
    ('⚖  法令遵守',      C_TEAL,   C_LTEAL,  '医療に関わる主要法令'),
]
for i, (label, dark, light, sub) in enumerate(badge_data):
    y = Inches(1.4 + i * 1.7)
    rect(s, Inches(5.8), y, Inches(7.0), Inches(1.5),
         fill=light, line=dark, lw=Pt(2))
    rect(s, Inches(5.8), y, Inches(0.18), Inches(1.5), fill=dark)
    tx(s, label,
       Inches(6.2), y + Inches(0.15), Inches(6.4), Inches(0.55),
       size=18, bold=True, color=dark)
    tx(s, sub,
       Inches(6.2), y + Inches(0.78), Inches(6.4), Inches(0.4),
       size=11, color=C_MGRAY)

# イラストプレースホルダー
rect(s, Inches(5.8), Inches(6.6), Inches(7.0), Inches(0.7),
     fill=RGBColor(0xE8, 0xF0, 0xFE), line=C_MGRAY)
tx(s, '★ ナノバナナ：cover_illust.png を挿入',
   Inches(6.0), Inches(6.72), Inches(6.6), Inches(0.35),
   size=9, color=C_MGRAY, align=PP_ALIGN.CENTER)

slide_no(s, 1)


# ============================================================
# スライド 2: アジェンダ
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_DARK, 'AGENDA', '本日のアジェンダ')

agenda = [
    ('PART 01', '個人情報保護',   '患者情報の正しい扱い方・漏洩防止',  C_BLUE,   C_LBLUE),
    ('PART 02', 'ハラスメント予防', '職場・対患者 ／ カスハラ重点解説', C_ORANGE, C_LORANGE),
    ('PART 03', '法令遵守',       '医療現場に関わる主要法令と義務',    C_TEAL,   C_LTEAL),
    ('WRAP UP', 'セルフチェック',  '自己確認チェックリスト',            C_DARK,   C_LGRAY),
]
xs = [Inches(0.5), Inches(6.9)]
ys = [Inches(1.1), Inches(4.2)]
positions = [(0,0),(0,1),(1,0),(1,1)]

for (row, col), (lbl, title, desc, dark, light) in zip(positions, agenda):
    x, y = xs[col], ys[row]
    rect(s, x, y, Inches(6.0), Inches(2.8), fill=light)
    rect(s, x, y, Inches(0.22), Inches(2.8), fill=dark)
    tx(s, lbl,
       x + Inches(0.35), y + Inches(0.2), Inches(5.4), Inches(0.35),
       size=9, bold=True, color=dark)
    tx(s, title,
       x + Inches(0.35), y + Inches(0.6), Inches(5.4), Inches(0.8),
       size=22, bold=True, color=C_DARK)
    tx(s, desc,
       x + Inches(0.35), y + Inches(1.55), Inches(5.4), Inches(0.4),
       size=12, color=C_MGRAY)

slide_no(s, 2)


# ============================================================
# スライド 3: 個人情報 — 定義
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_BLUE, 'PERSONAL INFORMATION  01', '個人情報保護　―　「個人情報」とは何か？')

# 左カード（薄ブルー）
rect(s, Inches(0.5), Inches(1.1), Inches(5.9), Inches(5.9), fill=C_LBLUE)
rect(s, Inches(0.5), Inches(1.1), Inches(5.9), Inches(0.55), fill=C_BLUE)
tx(s, '個人情報の定義', Inches(0.7), Inches(1.15), Inches(5.5), Inches(0.45),
   size=13, bold=True, color=C_WHITE)
items_l = ['氏名・生年月日・住所・電話番号',
           '診断名・病歴・検査結果・投薬内容',
           '顔写真・音声・映像',
           'マイナンバー・保険証番号']
for i, item in enumerate(items_l):
    rect(s, Inches(0.7), Inches(1.9 + i * 1.1), Inches(0.32), Inches(0.32),
         fill=C_BLUE)
    tx(s, item, Inches(1.15), Inches(1.87 + i * 1.1), Inches(5.0), Inches(0.5),
       size=13, color=C_DARK)

# 右カード（ダークブルー）
rect(s, Inches(6.9), Inches(1.1), Inches(5.9), Inches(5.9), fill=C_BLUE)
rect(s, Inches(6.9), Inches(1.1), Inches(5.9), Inches(0.55), fill=C_GOLD)
tx(s, '要配慮個人情報（特に厳格）', Inches(7.1), Inches(1.15), Inches(5.5), Inches(0.45),
   size=13, bold=True, color=C_DARK)
tx(s, '漏れただけで患者の人生に深刻な影響',
   Inches(7.1), Inches(1.8), Inches(5.5), Inches(0.4),
   size=11, color=C_MGRAY)
items_r = ['傷病・障害・健康診断結果',
           '精神疾患・感染症の罹患情報',
           '手術・検査・処方の記録']
for i, item in enumerate(items_r):
    rect(s, Inches(7.1), Inches(2.4 + i * 1.2), Inches(0.32), Inches(0.32),
         fill=C_GOLD)
    tx(s, item, Inches(7.55), Inches(2.37 + i * 1.2), Inches(5.0), Inches(0.5),
       size=14, bold=True, color=C_WHITE)

slide_no(s, 3)


# ============================================================
# スライド 4: 個人情報 — 5原則
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_BLUE, 'PERSONAL INFORMATION  01', '個人情報保護　―　患者情報取り扱い 5原則')

principles = [
    ('1', '目的外\n使用禁止',    '診療目的のみに使用'),
    ('2', '最小限の\nアクセス',   '担当外カルテは見ない'),
    ('3', '第三者提供\nの制限',   '本人同意なく渡さない'),
    ('4', '安全管理\n措置',       '施錠・ロックを徹底'),
    ('5', '漏洩時の\n即時報告',   '即座に上長へ連絡'),
]
for i, (num, label, hint) in enumerate(principles):
    x = Inches(0.5 + i * 2.58)
    # 番号円風（四角で代用）
    rect(s, x, Inches(1.1), Inches(2.35), Inches(1.5), fill=C_BLUE)
    tx(s, num, x, Inches(1.15), Inches(2.35), Inches(1.4),
       size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.6), Inches(2.35), Inches(1.55), fill=C_LBLUE)
    tx(s, label, x, Inches(2.65), Inches(2.35), Inches(1.4),
       size=14, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(4.15), Inches(2.35), Inches(0.8), fill=C_LGRAY)
    tx(s, hint, x, Inches(4.2), Inches(2.35), Inches(0.7),
       size=11, color=C_MGRAY, align=PP_ALIGN.CENTER)

# 矢印
for i in range(4):
    tx(s, '→', Inches(2.72 + i * 2.58), Inches(1.65), Inches(0.3), Inches(0.5),
       size=18, color=C_MGRAY, align=PP_ALIGN.CENTER)

footer_bar(s, C_BLUE, '「担当外の患者カルテを見る」だけで違反になります')
slide_no(s, 4)


# ============================================================
# スライド 5: 個人情報 — NG10
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_BLUE, 'PERSONAL INFORMATION  01', '個人情報保護　―　やってはいけないこと NG10')

ng_all = [
    ('院外でのカルテ・書類の閲覧・持ち出し',    True),
    ('私用スマホで患者情報を撮影・保存',          True),
    ('SNS への患者関連投稿（特定不能でも不可）', True),
    ('廊下・エレベーターでの患者情報の会話',      False),
    ('無関係な患者のカルテを閲覧する',            False),
    ('書類の放置・施錠忘れ',                      False),
    ('患者情報を含むメールを誤送信',              False),
    ('家族・知人への患者情報の口頭伝達',          False),
    ('退職後の情報の持ち出し・利用',              True),
    ('業務用 PC の共有パスワード使用',            False),
]
cols_x = [Inches(0.5), Inches(7.0)]
for idx, (item, highlight) in enumerate(ng_all):
    col = idx % 2
    row = idx // 2
    x = cols_x[col]
    y = Inches(1.15 + row * 1.3)
    bg = C_LRED if highlight else C_LGRAY
    rect(s, x, y, Inches(6.1), Inches(1.05), fill=bg)
    rect(s, x, y, Inches(0.55), Inches(1.05), fill=C_RED)
    tx(s, '✕', x, y + Inches(0.22), Inches(0.55), Inches(0.6),
       size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tx(s, item, x + Inches(0.65), y + Inches(0.28),
       Inches(5.3), Inches(0.5), size=12, color=C_DARK)

footer_bar(s, C_RED, '違反した場合 → 刑事罰（1年以下懲役・50万円以下罰金）・行政指導・民事賠償・社会的信用失墜')
slide_no(s, 5)


# ============================================================
# スライド 6: ハラスメント種類
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_ORANGE, 'HARASSMENT PREVENTION  02', 'ハラスメント予防　―　主な種類')

types = [
    ('パワハラ', C_DARK,   C_LGRAY,  '優越的立場を利用した言動で就業環境を害する',     '「なんでこんなこともできないんだ」と怒鳴る'),
    ('セクハラ', C_DARK,   C_LGRAY,  '性的言動により就業環境を害する',                 '体型や外見についての不適切な発言'),
    ('マタハラ', C_DARK,   C_LGRAY,  '妊娠・出産・育休に関する不利益な取り扱い',       '「妊娠するタイミングが悪い」などの発言'),
    ('カスハラ', C_RED,    C_LRED,   '患者・家族等による著しく不当な要求・言動',        '長時間怒鳴り・土下座要求・SNS脅迫　← 次スライドで詳解'),
]
for i, (name, dark, light, defn, ex) in enumerate(types):
    y = Inches(1.1 + i * 1.45)
    rect(s, Inches(0.5), y, Inches(2.3), Inches(1.25), fill=dark)
    tx(s, name, Inches(0.5), y + Inches(0.3), Inches(2.3), Inches(0.65),
       size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.9), y, Inches(9.9), Inches(1.25), fill=light)
    tx(s, defn, Inches(3.1), y + Inches(0.08), Inches(9.5), Inches(0.5),
       size=13, bold=True, color=dark)
    tx(s, f'例：{ex}', Inches(3.1), y + Inches(0.62), Inches(9.5), Inches(0.5),
       size=11.5, color=C_MGRAY)

rect(s, Inches(2.9), Inches(5.5 + 1.45), Inches(0.0), Inches(0.0))  # dummy
# カスハラ強調
tx(s, '★ 本日の重点テーマ',
   Inches(9.5), Inches(1.1 + 3 * 1.45) + Inches(0.35),
   Inches(2.8), Inches(0.55),
   size=12, bold=True, color=C_RED, align=PP_ALIGN.RIGHT)

slide_no(s, 6)


# ============================================================
# スライド 7: カスハラ — 定義 + 具体例
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_RED, 'CUSTOMER HARASSMENT  02 — 重点', 'カスタマーハラスメント（カスハラ）重点解説')

# 定義ボックス
rect(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.7), fill=C_RED)
tx(s, '定義（厚生労働省 2024年指針）',
   Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.38),
   size=9, color=RGBColor(0xFF, 0xCC, 0xCC))
tx(s, '顧客・患者・その家族等が行う、社会通念上不相当な言動で、\n労働者の就業環境が害されるもの',
   Inches(0.7), Inches(1.45), Inches(12.0), Inches(1.15),
   size=18, bold=True, color=C_WHITE)

# イラストプレースホルダー
rect(s, Inches(0.5), Inches(2.9), Inches(3.5), Inches(1.5),
     fill=C_LRED, line=C_RED, lw=Pt(1))
tx(s, '★ ナノバナナ：illust_kushara.png\n（怒る患者 vs 盾を持つスタッフ）',
   Inches(0.6), Inches(3.2), Inches(3.3), Inches(0.9),
   size=9, color=C_RED, align=PP_ALIGN.CENTER)

examples_4 = [
    ('① 暴言・脅迫系',   C_RED,    C_LRED,    '「訴えるぞ」「SNSで晒す」脅し\n大声の怒鳴り・罵倒・「上を呼べ」連続'),
    ('② 不当要求系',    C_ORANGE, C_LORANGE, '土下座・謝罪文の強要\n医学的に不要な処置・特別扱いの要求'),
    ('③ 身体的行為系',  C_DARK,   C_LGRAY,   '物を投げる・叩く\nスタッフへの不審な接触・追跡'),
    ('④ 長時間拘束系',  C_DARK,   C_LGRAY,   '何時間も居座りクレームを繰り返す\n電話での長時間クレーム'),
]
expos = [(0, 0), (0, 1), (1, 0), (1, 1)]
for (row, col), (title, dark, light, desc) in zip(expos, examples_4):
    x = Inches(4.2 + col * 4.6)
    y = Inches(2.9 + row * 2.2)
    rect(s, x, y, Inches(4.3), Inches(1.95), fill=light)
    rect(s, x, y, Inches(4.3), Inches(0.5), fill=dark)
    tx(s, title, x + Inches(0.12), y + Inches(0.08), Inches(4.0), Inches(0.38),
       size=12, bold=True, color=C_WHITE)
    tx(s, desc, x + Inches(0.12), y + Inches(0.6), Inches(4.0), Inches(1.25),
       size=12, color=C_DARK)

slide_no(s, 7)


# ============================================================
# スライド 8: カスハラ — 正当クレームとの違い
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_RED, 'CUSTOMER HARASSMENT  02', 'カスハラ　―　正当なクレームとの違い')

# 左：正当クレーム（グリーン）
rect(s, Inches(0.5), Inches(1.0), Inches(5.8), Inches(6.2), fill=RGBColor(0xE8,0xF8,0xED))
rect(s, Inches(0.5), Inches(1.0), Inches(5.8), Inches(0.65), fill=C_GREEN)
tx(s, '✓  正当なクレーム',
   Inches(0.7), Inches(1.08), Inches(5.4), Inches(0.5),
   size=17, bold=True, color=C_WHITE)
ok_items = [
    '待ち時間・説明不足への不満の申し出',
    '診療方針についての疑問・質問',
    'スタッフの対応への率直なフィードバック',
    '改善を求める具体的な要望',
]
for i, item in enumerate(ok_items):
    rect(s, Inches(0.7), Inches(1.85 + i * 0.9), Inches(0.35), Inches(0.35), fill=C_GREEN)
    tx(s, item, Inches(1.15), Inches(1.82 + i * 0.9), Inches(4.9), Inches(0.5),
       size=13, color=C_DARK)

rect(s, Inches(0.5), Inches(5.8), Inches(5.8), Inches(1.0), fill=C_GREEN)
tx(s, '→ 丁寧に受け止め、真摯に対応する',
   Inches(0.7), Inches(5.97), Inches(5.4), Inches(0.65),
   size=14, bold=True, color=C_WHITE)

# VS
tx(s, 'VS', Inches(6.35), Inches(3.6), Inches(0.6), Inches(0.55),
   size=22, bold=True, color=C_MGRAY, align=PP_ALIGN.CENTER)

# 右：カスハラ（レッド）
rect(s, Inches(7.0), Inches(1.0), Inches(5.8), Inches(6.2), fill=C_LRED)
rect(s, Inches(7.0), Inches(1.0), Inches(5.8), Inches(0.65), fill=C_RED)
tx(s, '✕  カスハラ',
   Inches(7.2), Inches(1.08), Inches(5.4), Inches(0.5),
   size=17, bold=True, color=C_WHITE)
ng_items_2 = [
    '怒鳴る・罵倒する・脅迫する',
    '土下座・謝罪文を強要する',
    '暴力行為・物を投げる',
    '長時間拘束・居座り・電話攻め',
]
for i, item in enumerate(ng_items_2):
    rect(s, Inches(7.2), Inches(1.85 + i * 0.9), Inches(0.35), Inches(0.35), fill=C_RED)
    tx(s, item, Inches(7.65), Inches(1.82 + i * 0.9), Inches(4.9), Inches(0.5),
       size=13, color=C_DARK)

rect(s, Inches(7.0), Inches(5.8), Inches(5.8), Inches(1.0), fill=C_RED)
tx(s, '→ 受け入れない。組織として対応する',
   Inches(7.2), Inches(5.97), Inches(5.4), Inches(0.65),
   size=14, bold=True, color=C_WHITE)

slide_no(s, 8)


# ============================================================
# スライド 9: カスハラ — 対応フロー
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_RED, 'CUSTOMER HARASSMENT  02', 'カスハラ　―　対応フロー 5ステップ')

steps = [
    ('STEP 1', '不当な言動に気づく',    '「これはクレームか、カスハラか？」を判断する',     C_DARK,   C_LGRAY),
    ('STEP 2', '一人で抱え込まない',    'すぐに同僚・先輩・リーダーへ声をかける　★最重要', C_ORANGE, C_LORANGE),
    ('STEP 3', '複数人で対応する',      '必ず2名以上。単独対応しない',                      C_DARK,   C_LGRAY),
    ('STEP 4', '事実を記録する',        '日時・言動の内容を具体的にメモする',               C_DARK,   C_LGRAY),
    ('STEP 5', '上長・窓口へ報告する',  '組織として対応。隠さない・一人で抱え込まない',     C_RED,    C_LRED),
]
for i, (num, title, desc, dark, light) in enumerate(steps):
    y = Inches(1.0 + i * 1.18)
    rect(s, Inches(0.5), y, Inches(1.6), Inches(1.0), fill=dark)
    tx(s, num, Inches(0.5), y + Inches(0.18), Inches(1.6), Inches(0.65),
       size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.25), y, Inches(10.55), Inches(1.0), fill=light)
    tx(s, title, Inches(2.45), y + Inches(0.05), Inches(3.8), Inches(0.5),
       size=14, bold=True, color=dark)
    tx(s, desc, Inches(6.3), y + Inches(0.08), Inches(6.3), Inches(0.85),
       size=13, color=C_DARK)
    if i < len(steps) - 1:
        tx(s, '▼', Inches(0.9), y + Inches(1.0), Inches(0.8), Inches(0.2),
           size=8, color=C_MGRAY, align=PP_ALIGN.CENTER)

footer_bar(s, C_RED,
    '【法的義務】2024年改正：事業者はカスハラへの対策・相談体制・スタッフフォローが義務化')
slide_no(s, 9)


# ============================================================
# スライド 10: 法令遵守
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_TEAL, 'COMPLIANCE  03', '法令遵守（コンプライアンス）　―　主要法令')

# イラストプレースホルダー
rect(s, Inches(10.8), Inches(1.1), Inches(2.0), Inches(1.6),
     fill=C_LTEAL, line=C_TEAL, lw=Pt(1))
tx(s, '★ ナノバナナ：\nillust_compliance.png',
   Inches(10.85), Inches(1.4), Inches(1.9), Inches(0.9),
   size=8, color=C_TEAL, align=PP_ALIGN.CENTER)

laws = [
    ('個人情報保護法',                  '患者情報の目的外利用・第三者提供禁止。\n漏洩時は72時間以内に報告義務'),
    ('医療法',                          '医療記録の正確な保管義務・\nインフォームドコンセントの確保'),
    ('労働施策総合推進法\n（パワハラ防止法）', '職場のパワハラ・カスハラ対策が事業者の義務。\n相談窓口の設置が必須'),
    ('医師法・保健師助産師看護師法',    '守秘義務（業務で知った情報を漏らしてはならない）。\n退職後も継続'),
    ('不正競争防止法',                  '院内の業務秘密の持ち出し・漏洩の禁止'),
]
for i, (law, point) in enumerate(laws):
    y = Inches(1.1 + i * 1.15)
    bg = C_TEAL if i % 2 == 0 else RGBColor(0x0A, 0x5C, 0x60)
    rect(s, Inches(0.5), y, Inches(3.0), Inches(1.05), fill=bg)
    tx(s, law, Inches(0.65), y + Inches(0.08), Inches(2.7), Inches(0.9),
       size=11, bold=True, color=C_WHITE)
    rect(s, Inches(3.6), y, Inches(7.1), Inches(1.05), fill=C_LTEAL)
    tx(s, point, Inches(3.75), y + Inches(0.08), Inches(6.8), Inches(0.9),
       size=11.5, color=C_DARK)

footer_bar(s, C_TEAL,
    '違反 → 【刑事罰】懲役・罰金　【行政処分】免許停止・取り消し　【民事賠償】損害賠償請求')
slide_no(s, 10)


# ============================================================
# スライド 11: 相談・報告
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_TEAL, 'COMPLIANCE  03', '法令遵守　―　迷ったら「報告・相談」が最優先')

action_steps = [
    (C_TEAL,   '「これはまずいかも？」と思ったら、まず直属の上長へ報告する'),
    (C_ORANGE, '上長への相談が難しい場合は院内コンプライアンス相談窓口を利用する'),
    (C_RED,    '隠さない・後回しにしない —— 早期対応が被害を最小化する'),
]
for i, (dark, text) in enumerate(action_steps):
    y = Inches(1.2 + i * 1.6)
    rect(s, Inches(0.5), y, Inches(0.9), Inches(1.3), fill=dark)
    tx(s, str(i + 1), Inches(0.5), y + Inches(0.25), Inches(0.9), Inches(0.8),
       size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(1.5), y, Inches(11.3), Inches(1.3), fill=C_LGRAY)
    rect(s, Inches(1.5), y, Inches(0.18), Inches(1.3), fill=dark)
    tx(s, text, Inches(1.8), y + Inches(0.38), Inches(10.8), Inches(0.55),
       size=15, color=C_DARK)

# 相談窓口
rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.3), fill=C_LTEAL)
rect(s, Inches(0.5), Inches(6.0), Inches(0.18), Inches(1.3), fill=C_TEAL)
tx(s, '院内相談窓口',
   Inches(0.8), Inches(6.05), Inches(12.0), Inches(0.45),
   size=13, bold=True, color=C_TEAL)
tx(s,
    '個人情報・コンプライアンス　担当部署：_______________　内線：___________\n'
    'ハラスメント　　　　　　　　担当部署：_______________　内線：___________',
    Inches(0.8), Inches(6.52), Inches(12.0), Inches(0.75),
    size=12, color=C_DARK)

slide_no(s, 11)


# ============================================================
# スライド 12: まとめ
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_DARK, '', 'まとめ　― 3つのメッセージ')

messages = [
    (C_BLUE,   '01', '患者情報はあなたが「預かっている」もの',
                     '慎重に、目的の範囲内だけで扱ってください。'),
    (C_RED,    '02', 'ハラスメントは職場の問題',
                     '特にカスハラは一人で抱え込まず、必ず組織で対応してください。'),
    (C_TEAL,   '03', '迷ったら報告',
                     '隠さないことが、自分と患者と病院を守ります。'),
]
for i, (dark, num, title, desc) in enumerate(messages):
    y = Inches(1.1 + i * 1.85)
    light = C_LBLUE if dark == C_BLUE else (C_LRED if dark == C_RED else C_LTEAL)
    rect(s, Inches(0.5), y, Inches(1.4), Inches(1.65), fill=dark)
    tx(s, num, Inches(0.5), y + Inches(0.35), Inches(1.4), Inches(0.95),
       size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.0), y, Inches(11.0), Inches(1.65), fill=light)
    rect(s, Inches(2.0), y, Inches(0.18), Inches(1.65), fill=dark)
    tx(s, title, Inches(2.35), y + Inches(0.15), Inches(10.5), Inches(0.6),
       size=18, bold=True, color=dark)
    tx(s, desc, Inches(2.35), y + Inches(0.85), Inches(10.5), Inches(0.65),
       size=14, color=C_DARK)

footer_bar(s, C_DARK,
    '研修後はセルフチェックリスト（配布資料 最終ページ）で自己確認してください')
slide_no(s, 12)


# ============================================================
# スライド 13: セルフチェックリスト
# ============================================================
s = prs.slides.add_slide(BLANK)
header_bar(s, C_GOLD, 'SELF CHECK', 'セルフチェックリスト　― 今日から使える')

checks = [
    (C_BLUE,   '個人情報保護', [
        '患者情報を院外・廊下で話していない',
        '私用スマホで患者情報を撮影していない',
        'SNS への患者関連投稿を一切していない',
        '業務終了時に書類を施錠保管している',
    ]),
    (C_RED,    'ハラスメント予防', [
        '指導の際に怒鳴る・人格否定をしていない',
        '患者から不当な言動を受けたとき報告している',
        'カスハラと正当クレームを区別できる',
        'カスハラ対応は複数人で行っている',
    ]),
    (C_TEAL,   '法令遵守', [
        '守秘義務（退職後も含む）を認識している',
        '迷ったときは上長・相談窓口に報告できている',
        'インシデントは速やかに報告している',
    ]),
]
xs = [Inches(0.5), Inches(4.7), Inches(9.0)]
for j, (dark, cat, items) in enumerate(checks):
    light = C_LBLUE if dark == C_BLUE else (C_LRED if dark == C_RED else C_LTEAL)
    x = xs[j]
    w = Inches(3.9)
    rect(s, x, Inches(1.05), w, Inches(0.55), fill=dark)
    tx(s, cat, x + Inches(0.1), Inches(1.1), w - Inches(0.2), Inches(0.45),
       size=12, bold=True, color=C_WHITE)
    for k, item in enumerate(items):
        y_c = Inches(1.75 + k * 1.3)
        rect(s, x, y_c, w, Inches(1.1), fill=light)
        rect(s, x + Inches(0.18), y_c + Inches(0.35),
             Inches(0.38), Inches(0.38),
             fill=C_WHITE, line=dark, lw=Pt(1.5))
        tx(s, item,
           x + Inches(0.7), y_c + Inches(0.25),
           w - Inches(0.85), Inches(0.65),
           size=11.5, color=C_DARK)

slide_no(s, 13)


# ============================================================
# 保存
# ============================================================
out = 'C:/Users/yuuna/agens/勉強会/研修スライド_カラー版.pptx'
prs.save(out)
print(f'保存完了：{out}')
print(f'スライド数：{len(prs.slides)}枚')
