# -*- coding: utf-8 -*-
"""新人勉強会 - 移乗・移動介助 研修資料生成スクリプト"""

from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import copy

# ─────────────────────────────────────
#  カラー定数
# ─────────────────────────────────────
C_PRIMARY   = RGBColor(0x1A, 0x56, 0xDB)   # 濃いブルー
C_ACCENT    = RGBColor(0x0E, 0x9F, 0x6E)   # グリーン
C_WARNING   = RGBColor(0xE3, 0xA0, 0x08)   # アンバー
C_DANGER    = RGBColor(0xE0, 0x28, 0x28)   # レッド
C_LIGHT_BG  = RGBColor(0xEF, 0xF6, 0xFF)   # 薄いブルー背景
C_DARK      = RGBColor(0x1F, 0x29, 0x37)   # ほぼ黒
C_GRAY      = RGBColor(0x6B, 0x72, 0x80)   # グレー
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

def set_cell_bg(cell, r, g, b):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), f'{r:02X}{g:02X}{b:02X}')
    tcPr.append(shd)

def set_para_bg(para, r, g, b):
    pPr = para._p.get_or_add_pPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), f'{r:02X}{g:02X}{b:02X}')
    pPr.append(shd)

def add_colored_heading(doc, text, level=1, color=C_PRIMARY, size=16, bold=True, bg=None):
    p = doc.add_paragraph()
    if bg:
        set_para_bg(p, *bg)
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return p

def add_step_table(doc, steps, title="手順"):
    """番号付き手順テーブル"""
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    # ヘッダー
    hdr = table.rows[0].cells
    for i, txt in enumerate([" No.", f" {title}", " ポイント"]):
        hdr[i].text = txt
        hdr[i].paragraphs[0].runs[0].bold = True
        hdr[i].paragraphs[0].runs[0].font.color.rgb = C_WHITE
        hdr[i].paragraphs[0].runs[0].font.size = Pt(9)
        set_cell_bg(hdr[i], 0x1A, 0x56, 0xDB)
    # 列幅
    for i, w in enumerate([Cm(1.2), Cm(7.5), Cm(6.3)]):
        for cell in table.columns[i].cells:
            cell.width = w
    for step in steps:
        row = table.add_row().cells
        row[0].text = str(step[0])
        row[0].paragraphs[0].runs[0].font.size = Pt(10)
        row[0].paragraphs[0].runs[0].bold = True
        row[0].paragraphs[0].runs[0].font.color.rgb = C_PRIMARY
        row[1].text = step[1]
        row[1].paragraphs[0].runs[0].font.size = Pt(9)
        row[2].text = step[2]
        row[2].paragraphs[0].runs[0].font.size = Pt(9)
        row[2].paragraphs[0].runs[0].font.color.rgb = C_GRAY
        if step[0] % 2 == 0:
            set_cell_bg(row[0], 0xEB, 0xF5, 0xFF)
            set_cell_bg(row[1], 0xEB, 0xF5, 0xFF)
            set_cell_bg(row[2], 0xEB, 0xF5, 0xFF)
    return table

def add_check_box(doc, items, title="チェックリスト"):
    p = doc.add_paragraph()
    run = p.add_run(f"☑ {title}")
    run.bold = True
    run.font.size = Pt(10)
    run.font.color.rgb = C_ACCENT
    for item in items:
        p2 = doc.add_paragraph(style='List Bullet')
        p2.paragraph_format.left_indent = Cm(0.5)
        run2 = p2.add_run(f"□ {item}")
        run2.font.size = Pt(9)

def add_warning_box(doc, text, color=(0xFF, 0xF3, 0xCD), icon="⚠️"):
    table = doc.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    cell.text = f"{icon}  {text}"
    cell.paragraphs[0].runs[0].font.size = Pt(9)
    cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x92, 0x40, 0x09)
    set_cell_bg(cell, *color)
    return table

def add_image_prompt_box(doc, slide_no, prompt_ja, prompt_en):
    """画像生成プロンプトボックス（発表原稿用）"""
    table = doc.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    p = cell.paragraphs[0]
    run = p.add_run(f"🎨 [スライド{slide_no} 画像生成プロンプト]\n")
    run.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x4A, 0x1D, 0x96)
    run2 = p.add_run(f"日本語: {prompt_ja}\n")
    run2.font.size = Pt(9)
    run2.font.color.rgb = C_DARK
    run3 = p.add_run(f"English: {prompt_en}")
    run3.font.italic = True
    run3.font.size = Pt(8.5)
    run3.font.color.rgb = C_GRAY
    set_cell_bg(cell, 0xED, 0xEB, 0xFF)
    return table


# ═══════════════════════════════════════
#  1. 配布資料（Word）
# ═══════════════════════════════════════
def create_handout():
    doc = Document()

    # ページ余白
    for sec in doc.sections:
        sec.top_margin    = Cm(2.0)
        sec.bottom_margin = Cm(2.0)
        sec.left_margin   = Cm(2.2)
        sec.right_margin  = Cm(2.2)

    # ─── 表紙ヘッダー ───────────────────────────
    p_title = doc.add_paragraph()
    set_para_bg(p_title, 0x1A, 0x56, 0xDB)
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_title.paragraph_format.space_before = Pt(8)
    p_title.paragraph_format.space_after  = Pt(8)
    r1 = p_title.add_run("新人研修テキスト｜看護・介護職員向け")
    r1.font.size = Pt(10); r1.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF); r1.bold = True

    p_main = doc.add_paragraph()
    set_para_bg(p_main, 0x1A, 0x56, 0xDB)
    p_main.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_main.paragraph_format.space_after = Pt(4)
    rm = p_main.add_run("移乗・移動介助の基本")
    rm.font.size = Pt(26); rm.font.color.rgb = C_WHITE; rm.bold = True

    p_sub = doc.add_paragraph()
    set_para_bg(p_sub, 0x1A, 0x56, 0xDB)
    p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_sub.paragraph_format.space_after = Pt(10)
    rs = p_sub.add_run("Transfer & Mobility Assistance  ｜  所要時間：約60分")
    rs.font.size = Pt(10); rs.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)

    # サブ情報行
    info_tbl = doc.add_table(rows=1, cols=3)
    info_tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    for cell, txt in zip(info_tbl.rows[0].cells,
                         ["📅 実施日：　　　　", "👤 担当講師：　　　　", "✏️ 受講者名：　　　　"]):
        cell.text = txt
        cell.paragraphs[0].runs[0].font.size = Pt(9)
        set_cell_bg(cell, 0xEB, 0xF5, 0xFF)
    doc.add_paragraph()

    # ─── 研修のねらい ────────────────────────────
    add_colored_heading(doc, "📋 研修のねらい", size=13, color=C_PRIMARY)
    aims = [
        "ボディメカニクスの原則を理解し、自分の体を守りながら安全に介助できる",
        "ベッド↔車椅子、車椅子↔トイレの移乗手順を正確に覚える",
        "歩行・車椅子移動の介助を安全に実践できる",
        "転倒・転落リスクを事前に察知し、予防行動がとれる",
        "利用者の「残存能力」を活かした自立支援の視点を持つ",
    ]
    for a in aims:
        p = doc.add_paragraph(style='List Bullet')
        p.paragraph_format.left_indent = Cm(0.5)
        r = p.add_run(a)
        r.font.size = Pt(9.5)

    doc.add_paragraph()

    # ─── 目次 ────────────────────────────────────
    add_colored_heading(doc, "目  次", size=12, color=C_PRIMARY)
    toc = [
        ("1", "介助の3大原則", "5分", "p.2"),
        ("2", "ボディメカニクスとは", "10分", "p.3"),
        ("3", "移乗介助①　ベッド↔車椅子", "15分", "p.5"),
        ("4", "移乗介助②　車椅子↔トイレ", "10分", "p.9"),
        ("5", "移動介助　歩行・車椅子", "12分", "p.12"),
        ("6", "安全管理とリスク予防", "6分", "p.15"),
        ("7", "チェックリスト・まとめ", "2分", "p.17"),
    ]
    toc_tbl = doc.add_table(rows=1, cols=4)
    toc_tbl.style = 'Table Grid'
    for cell, hdr in zip(toc_tbl.rows[0].cells, [" #", " セクション", " 時間", " ページ"]):
        cell.text = hdr
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(9)
        cell.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        set_cell_bg(cell, 0x1A, 0x56, 0xDB)
    for row_data in toc:
        r = toc_tbl.add_row().cells
        for i, v in enumerate(row_data):
            r[i].text = f" {v}"
            r[i].paragraphs[0].runs[0].font.size = Pt(9)
    doc.add_page_break()

    # ─── SECTION 1: 介助の3大原則 ────────────────
    add_colored_heading(doc, "1. 介助の3大原則", size=14, color=C_PRIMARY)
    p = doc.add_paragraph()
    r = p.add_run("すべての介助に共通する根本的な考え方です。この3つを常に意識してください。")
    r.font.size = Pt(9.5)

    principles = [
        ("🛡️ 安全", (0xE0, 0x28, 0x28),
         "転倒・転落・皮膚損傷を防ぐ",
         "ブレーキ確認、環境整備、声かけを怠らない。「大丈夫だろう」は禁物。"),
        ("😌 安楽", (0x0E, 0x9F, 0x6E),
         "利用者に苦痛・不安を与えない",
         "急な動き・冷たい手・無言の介助はNG。丁寧な声かけとゆっくりした動作を心がける。"),
        ("🤝 自立支援", (0x1A, 0x56, 0xDB),
         "できることは本人にやってもらう",
         "全部やってしまわない。「少しだけ手伝う」「一緒にやる」スタンスが大切。"),
    ]
    for icon_title, color, short, detail in principles:
        tbl = doc.add_table(rows=1, cols=2)
        tbl.style = 'Table Grid'
        c0, c1 = tbl.rows[0].cells
        c0.width = Cm(3.5); c1.width = Cm(12)
        c0.text = icon_title
        c0.paragraphs[0].runs[0].bold = True
        c0.paragraphs[0].runs[0].font.size = Pt(11)
        c0.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        c0.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        set_cell_bg(c0, *color)
        p1 = c1.paragraphs[0]
        r1 = p1.add_run(short + "\n")
        r1.bold = True; r1.font.size = Pt(10)
        r2 = p1.add_run(detail)
        r2.font.size = Pt(9); r2.font.color.rgb = C_GRAY
        doc.add_paragraph()

    add_warning_box(doc, "「転ばせてしまった」は最大のインシデント。スピードより確実さを優先しましょう。")
    doc.add_paragraph()

    # ─── SECTION 2: ボディメカニクス ─────────────
    add_colored_heading(doc, "2. ボディメカニクスとは", size=14, color=C_PRIMARY)
    p = doc.add_paragraph()
    r = p.add_run(
        "「ボディメカニクス」とは、人体の構造や物理の法則を活かして、"
        "介助者の体への負担を最小限にしながら安全に介助するための技術です。\n"
        "腰痛は介護職の離職理由No.1。正しい姿勢を身につければ長く働けます。"
    )
    r.font.size = Pt(9.5)
    doc.add_paragraph()

    bm_steps = [
        (1, "足を肩幅に開く（支持基底面を広げる）",
         "足の間隔が広いほど安定する。ふらつき防止の第一歩。"),
        (2, "膝を曲げて重心を低くする",
         "腰ではなく脚の筋肉で支える。立ち上がるときも膝から。"),
        (3, "利用者にできるだけ近づく",
         "腕を伸ばして介助すると腰への負担が激増する。密着が基本。"),
        (4, "大きな筋肉（下肢・腹部）を使う",
         "背中の小さな筋肉ではなく、太ももや体幹で動かす。"),
        (5, "体をねじらない　→ 足ごと向きを変える",
         "ひねり動作が椎間板を傷める最大の原因。"),
        (6, "水平・てこの原理を活用する",
         "持ち上げるより引く・回す・滑らせる動作の方が楽。"),
    ]
    add_step_table(doc, bm_steps, title="ボディメカニクス 6原則")
    doc.add_paragraph()

    # NG例テーブル
    add_colored_heading(doc, "🚫 NG姿勢 vs ✅ OK姿勢", size=11, color=C_DANGER)
    ng_tbl = doc.add_table(rows=1, cols=2)
    ng_tbl.style = 'Table Grid'
    for cell, txt in zip(ng_tbl.rows[0].cells, ["🚫 やってしまいがちなNG", "✅ 正しい対応"]):
        cell.text = txt
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(9)
        cell.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        bg = (0xC8, 0x1E, 0x1E) if "NG" in txt else (0x05, 0x7A, 0x55)
        set_cell_bg(cell, *bg)
    ng_ok = [
        ("腕だけで引っ張り上げる", "体を密着させて体幹で支える"),
        ("前かがみで腰を曲げる", "膝を曲げて骨盤を低くする"),
        ("立ったまま上から押し下げる", "目線を合わせてから動作開始"),
        ("急いで素早く動かす", "「せーの」と声をかけてゆっくり"),
        ("一人で無理に抱える", "遠慮せずにスタッフを呼ぶ"),
    ]
    for ng, ok in ng_ok:
        row = ng_tbl.add_row().cells
        row[0].text = f"✕ {ng}"
        row[0].paragraphs[0].runs[0].font.size = Pt(9)
        row[1].text = f"○ {ok}"
        row[1].paragraphs[0].runs[0].font.size = Pt(9)
    doc.add_page_break()

    # ─── SECTION 3: 移乗介助① ベッド↔車椅子 ──────
    add_colored_heading(doc, "3. 移乗介助①　ベッド → 車椅子", size=14, color=C_PRIMARY)
    p = doc.add_paragraph()
    r = p.add_run(
        "移乗はベッドと車椅子を行き来する最も基本的な介助です。\n"
        "「準備8割、介助2割」—— 環境整備が最重要です。"
    )
    r.font.size = Pt(9.5)
    doc.add_paragraph()

    add_colored_heading(doc, "◆ 事前準備チェック", size=11, color=C_ACCENT)
    pre_checks = [
        "車椅子のブレーキがかかっているか確認",
        "フットサポート（足置き）を跳ね上げているか確認",
        "ベッドの高さを車椅子の座面に合わせているか（端座位で足底が床につく高さ）",
        "介助者自身の足元・周囲に障害物がないか",
        "利用者に介助内容を説明し、同意を得る",
        "滑り止めマット・スライディングボードなど補助具の準備",
    ]
    add_check_box(doc, pre_checks, "準備チェック（介助開始前に必ず確認）")
    doc.add_paragraph()

    steps_bed2wc = [
        (1, "ベッドの高さを調整し、端座位を取ってもらう",
         "両足が床につく高さ。利用者の横に立つ。"),
        (2, "車椅子をベッドに対して20〜30°の角度で置く",
         "健側（麻痺がない方）に車椅子を置く。"),
        (3, "利用者に前傾姿勢をとってもらう（鼻がひざの上に来るように）",
         "重心が前に移動→立ちやすくなる。"),
        (4, "介助者は利用者の患側（麻痺側）に立ち、腰に手を回す",
         "膝で利用者の膝をブロックするように立つ。"),
        (5, "「せーの」で一緒に立ち上がる（介助者は膝を伸ばして持ち上げない）",
         "利用者の動きに合わせて補助する意識で。"),
        (6, "立位が安定したら、健側方向に小さく方向転換する",
         "足をすり足で動かし、ひねらない。"),
        (7, "車椅子の座面を確認してから、ゆっくり着座させる",
         "深く座れているか確認。ずり落ち防止。"),
        (8, "フットサポートを戻し、姿勢を整えて声かけ",
         "「座れましたか？」「痛いところはないですか？」"),
    ]
    add_step_table(doc, steps_bed2wc, title="ベッド → 車椅子 移乗手順")
    doc.add_paragraph()

    add_colored_heading(doc, "◆ 車椅子 → ベッド（逆手順のポイント）", size=11, color=C_ACCENT)
    steps_wc2bed = [
        (1, "ベッドに対して健側から20〜30°の角度で車椅子を近づける",
         "患側から戻すと方向転換が大変になる。"),
        (2, "ブレーキをかけ、フットサポートを跳ね上げる",
         "絶対に確認。ブレーキ忘れは転倒に直結。"),
        (3, "前傾姿勢から立ち上がり、ベッドに向かって方向転換",
         "ベッド用の動作はほぼ逆手順。"),
        (4, "ゆっくりと着座し、仰臥位（横になる）への移行を手伝う",
         "足をベッドに乗せる際は支えながらゆっくりと。"),
    ]
    add_step_table(doc, steps_wc2bed, title="車椅子 → ベッド 移乗手順（要点）")
    doc.add_paragraph()

    add_warning_box(doc,
        "ブレーキかけ忘れは転倒事故の主原因です。移乗前の「ブレーキ確認」を声に出して習慣化しましょう。",
        icon="⚠️")
    doc.add_page_break()

    # ─── SECTION 4: 移乗介助② 車椅子↔トイレ ──────
    add_colored_heading(doc, "4. 移乗介助②　車椅子 → トイレ", size=14, color=C_PRIMARY)
    p = doc.add_paragraph()
    r = p.add_run(
        "トイレ介助は利用者の尊厳に深く関わります。\n"
        "プライバシーへの配慮と手際のよさの両立が求められます。"
    )
    r.font.size = Pt(9.5)
    doc.add_paragraph()

    steps_toilet = [
        (1, "トイレ内の安全確認（床の水濡れ・手すりの確認）",
         "濡れた床は転倒リスク大。必ず拭く。"),
        (2, "車椅子を便座に対して斜め45°に置く（健側手すり側）",
         "手すりをつかみやすい配置にする。"),
        (3, "ブレーキ確認・フットサポート跳ね上げ",
         "忘れず声に出してチェック。"),
        (4, "衣服を下ろし、立位介助",
         "できる限り利用者自身にやってもらう。"),
        (5, "手すりにつかまりながら方向転換 → 着座",
         "座面の位置を確認してからゆっくり。"),
        (6, "カーテンを引いてプライバシー確保／ナースコール位置を伝える",
         "「終わったら押してください」と必ず伝える。"),
        (7, "終了後：立ち上がり、清潔操作、衣服を整える",
         "できる動作は本人に。"),
        (8, "手洗いを介助し、車椅子へ移乗して退出",
         "手指衛生は感染予防の基本。"),
    ]
    add_step_table(doc, steps_toilet, title="車椅子 → トイレ 手順")
    doc.add_paragraph()

    add_colored_heading(doc, "◆ 尊厳への配慮", size=11, color=C_ACCENT)
    dignity = doc.add_paragraph()
    dignity_text = (
        "トイレは最もプライベートな行為です。\n"
        "• 無駄に中に入らない / 用が済んだらすぐ出る\n"
        "• 「〜してあげる」ではなく「一緒にやりましょう」の言葉を使う\n"
        "• 長い待ち時間はNG → スムーズに動ける準備を事前に整える\n"
        "• 失禁があっても表情・声で驚かない、責めない"
    )
    r = dignity.add_run(dignity_text)
    r.font.size = Pt(9.5)
    doc.add_page_break()

    # ─── SECTION 5: 移動介助 ──────────────────────
    add_colored_heading(doc, "5. 移動介助", size=14, color=C_PRIMARY)

    # 5-1 歩行介助
    add_colored_heading(doc, "◆ 5-1　歩行介助", size=12, color=C_ACCENT)
    p = doc.add_paragraph()
    r = p.add_run(
        "歩行介助は「転倒させない」より「自分で歩ける自信をつけてもらう」が目標です。"
    )
    r.font.size = Pt(9.5)
    doc.add_paragraph()

    walk_steps = [
        (1, "患側（麻痺・弱い側）の斜め後ろに立つ",
         "利用者が倒れかけたとき最も支えやすい位置。"),
        (2, "腰〜骨盤に手を添える（つかまない、支える）",
         "脇を抱えるのはNG→肩関節脱臼のリスク。"),
        (3, "利用者の歩幅・ペースに合わせる",
         "「急かさない」が大原則。焦りは転倒を招く。"),
        (4, "段差・障害物を先に伝える",
         "「少し段があります」と手前で声かけ。"),
        (5, "疲労感・歩き方の変化を観察しながら歩く",
         "足の引きずり増加→休憩のサイン。"),
    ]
    add_step_table(doc, walk_steps, title="歩行介助の手順")
    doc.add_paragraph()

    # 杖・歩行器
    add_colored_heading(doc, "◆ 補助具使用時のポイント", size=11, color=C_ACCENT)
    aids_tbl = doc.add_table(rows=1, cols=3)
    aids_tbl.style = 'Table Grid'
    for cell, txt in zip(aids_tbl.rows[0].cells, [" 補助具", " 上り（段差）", " 下り（段差）"]):
        cell.text = txt
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(9)
        cell.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        set_cell_bg(cell, 0x1A, 0x56, 0xDB)
    for aid, up, down in [
        ("杖", "杖→健足→患足（杖先行）", "杖→患足→健足"),
        ("歩行器", "歩行器を先に前進→両足を揃える", "歩行器を先に降ろす→両足"),
        ("手すり", "健手で手すりをつかむ", "患側から先に降りる"),
    ]:
        row = aids_tbl.add_row().cells
        for i, v in enumerate([aid, up, down]):
            row[i].text = f" {v}"
            row[i].paragraphs[0].runs[0].font.size = Pt(9)
    doc.add_paragraph()

    # 5-2 車椅子移動
    add_colored_heading(doc, "◆ 5-2　車椅子での移動介助", size=12, color=C_ACCENT)
    wc_steps = [
        (1, "出発前：ブレーキ解除・フットサポートに両足が乗っているか確認",
         "足が引きずられる事故を防ぐ。"),
        (2, "後ろから両手でグリップを持ち、視線を進行方向に向ける",
         "利用者の頭頂部より高い視点を保つ。"),
        (3, "段差（スロープ下り）：後ろ向きにゆっくり",
         "前向きに降りると利用者が前傾して危険。"),
        (4, "段差（スロープ上り）：前向きに角度を保って一定スピードで",
         "急停止はふらつきの原因。"),
        (5, "エレベーター：後ろ向きで乗り込み、前向きで降りる",
         "扉に挟まれないよう前方を確認。"),
        (6, "走行中：利用者への声かけを忘れずに",
         "「段差があります」「曲がります」など。"),
    ]
    add_step_table(doc, wc_steps, title="車椅子移動の手順・注意点")
    doc.add_page_break()

    # ─── SECTION 6: 安全管理・リスク予防 ────────────
    add_colored_heading(doc, "6. 安全管理とリスク予防", size=14, color=C_PRIMARY)

    risk_tbl = doc.add_table(rows=1, cols=3)
    risk_tbl.style = 'Table Grid'
    for cell, txt in zip(risk_tbl.rows[0].cells, [" リスク", " 原因（主な）", " 予防策"]):
        cell.text = txt
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(9)
        cell.paragraphs[0].runs[0].font.color.rgb = C_WHITE
        set_cell_bg(cell, 0xC8, 0x1E, 0x1E)
    risks = [
        ("転倒・転落", "ブレーキ忘れ、焦り、環境不備",
         "ブレーキ声出し確認、環境整備、二人介助"),
        ("腰痛（介助者）", "前かがみ、腕だけで支える",
         "ボディメカニクス徹底、補助具活用"),
        ("皮膚損傷", "ずり落ち、引きずり、摩擦",
         "水平移動、スライディングシート活用"),
        ("肩関節脱臼", "脇を持って引き上げる",
         "骨盤・体幹を支える、脇はNG"),
        ("誤嚥（移動後）", "食後すぐの移動・急な体位変換",
         "食後30分は安静、移動前に確認"),
    ]
    for r_data in risks:
        row = risk_tbl.add_row().cells
        for i, v in enumerate(r_data):
            row[i].text = f" {v}"
            row[i].paragraphs[0].runs[0].font.size = Pt(9)
    doc.add_paragraph()

    add_warning_box(doc,
        "ヒヤリハットは必ず報告！「大事にならなかったから報告しなくていい」は誤りです。"
        "記録と共有が次の事故を防ぎます。", icon="⚠️")
    doc.add_paragraph()

    # 観察ポイント
    add_colored_heading(doc, "◆ 介助中に観察すること", size=11, color=C_ACCENT)
    obs = [
        "顔色・表情（痛みの徴候：眉間のしわ、顔のゆがみ）",
        "息切れ・めまい・冷や汗（循環動態の変化）",
        "言語的コミュニケーション（「痛い」「だるい」「くらくらする」）",
        "皮膚（赤み・傷・発疹）",
        "歩行バランス・可動域の変化（昨日と違う）",
    ]
    add_check_box(doc, obs, "介助中の観察チェック")
    doc.add_page_break()

    # ─── SECTION 7: まとめ・チェックリスト ───────────
    add_colored_heading(doc, "7. まとめ＆自己チェックリスト", size=14, color=C_PRIMARY)

    summary_p = doc.add_paragraph()
    r = summary_p.add_run(
        "今日学んだことを振り返り、自分の理解度を確認しましょう。\n"
        "「わかった」から「できる」へ——繰り返し実践することが大切です。"
    )
    r.font.size = Pt(9.5)
    doc.add_paragraph()

    final_checks = [
        "ボディメカニクス6原則をすらすら言えますか？",
        "ベッド→車椅子の移乗手順を8ステップで説明できますか？",
        "ブレーキ確認を声に出す習慣がつきましたか？",
        "患側（麻痺側）に立って歩行介助できますか？",
        "段差での車椅子操作（後ろ向き下り）を理解しましたか？",
        "転倒リスク5つを挙げられますか？",
        "トイレ介助で尊厳への配慮ができていますか？",
        "ヒヤリハットは必ず報告するという意識が持てましたか？",
    ]
    add_check_box(doc, final_checks, "研修後 自己チェックリスト")
    doc.add_paragraph()

    # キーワードまとめ
    add_colored_heading(doc, "◆ 今日の重要キーワード", size=11, color=C_ACCENT)
    kw_tbl = doc.add_table(rows=3, cols=3)
    kw_tbl.style = 'Table Grid'
    keywords = [
        ("ボディメカニクス", "支持基底面", "重心を低く"),
        ("患側に立つ", "ブレーキ確認", "前傾姿勢"),
        ("自立支援", "ヒヤリハット報告", "安全・安楽・自立"),
    ]
    for i, row_kw in enumerate(keywords):
        for j, kw in enumerate(row_kw):
            cell = kw_tbl.cell(i, j)
            cell.text = f"  {kw}"
            cell.paragraphs[0].runs[0].font.size = Pt(9.5)
            cell.paragraphs[0].runs[0].bold = True
            cell.paragraphs[0].runs[0].font.color.rgb = C_PRIMARY
            if (i + j) % 2 == 0:
                set_cell_bg(cell, 0xEB, 0xF5, 0xFF)
    doc.add_paragraph()

    # フッター
    p_footer = doc.add_paragraph()
    set_para_bg(p_footer, 0x1A, 0x56, 0xDB)
    p_footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_f = p_footer.add_run("新人研修テキスト｜移乗・移動介助の基本　©研修事務局")
    r_f.font.size = Pt(8)
    r_f.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)

    doc.save("C:/Users/yuuna/agens/新人勉強会/配布資料_移乗移動介助.docx")
    print("OK: Word handout saved")


# ═══════════════════════════════════════
#  2. 発表原稿＋画像プロンプト（Word）
# ═══════════════════════════════════════
def create_script():
    doc = Document()
    for sec in doc.sections:
        sec.top_margin = Cm(2.0)
        sec.bottom_margin = Cm(2.0)
        sec.left_margin = Cm(2.5)
        sec.right_margin = Cm(2.5)

    # タイトル
    p = doc.add_paragraph()
    r = p.add_run("発表原稿 ＋ 画像生成プロンプト集")
    r.bold = True; r.font.size = Pt(18)
    r.font.color.rgb = C_PRIMARY
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2 = doc.add_paragraph()
    r2 = p2.add_run("移乗・移動介助 新人研修　PowerPointスライド用")
    r2.font.size = Pt(11); r2.font.color.rgb = C_GRAY
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    note_p = doc.add_paragraph()
    r_note = note_p.add_run(
        "【使い方】\n"
        "・各スライドに対応した発表原稿（話し言葉）と画像生成プロンプトが記載されています。\n"
        "・プロンプトはMidjourney / DALL·E 3 / Stable Diffusion 等で使用できます。\n"
        "・紫のボックスが画像プロンプト、下の本文が発表原稿です。"
    )
    r_note.font.size = Pt(9.5)
    doc.add_paragraph()

    slides = [
        {
            "no": 1,
            "title": "タイトルスライド",
            "time": "1分",
            "prompt_ja": "明るい介護施設の廊下で、若い看護師が笑顔でベッドから車椅子へ移乗介助をしている場面。柔らかい自然光、ウォームカラー、プロフェッショナルで清潔感のある雰囲気",
            "prompt_en": "A bright nursing facility hallway, young nurse smiling while assisting transfer from bed to wheelchair, soft natural lighting, warm colors, clean professional medical setting, photorealistic",
            "script": (
                "皆さん、本日はよろしくお願いします。\n"
                "今日は「移乗・移動介助の基本」について、約35分で学びます。\n\n"
                "移乗や移動介助は、毎日何十回も行う業務です。\n"
                "でも「なんとなくやっている」という状態が最も危険です。\n\n"
                "今日の研修で、なぜそうするのかを理解した上で動けるようになりましょう。"
            )
        },
        {
            "no": 2,
            "title": "今日のゴール（研修目標）",
            "time": "1.5分",
            "prompt_ja": "5つのアイコン（盾・人・手・歩行・目）がシンプルにレイアウトされた、フラットデザインのインフォグラフィック。青とグリーンのカラースキーム、クリーンな白背景",
            "prompt_en": "Flat design infographic with 5 icons representing safety, person, hands, walking, and eye. Blue and green color scheme, clean white background, modern minimal style",
            "script": (
                "今日の研修で身につけてほしいことは5つです。\n\n"
                "①ボディメカニクスの原則を理解すること。\n"
                "②ベッドと車椅子の移乗手順を正確に覚えること。\n"
                "③歩行介助を安全に行えること。\n"
                "④転倒リスクを予測して予防できること。\n"
                "⑤利用者の「できること」を大切にした自立支援の視点を持つこと。\n\n"
                "特に⑤は見落とされがちですが、介護・看護の根幹です。\n"
                "「やってあげる」ではなく「一緒にやる」という姿勢を大切にしてください。"
            )
        },
        {
            "no": 3,
            "title": "介助の3大原則",
            "time": "2分",
            "prompt_ja": "3つの大きなアイコン（盾・笑顔・握手）が横並びに並ぶ、シンプルなインフォグラフィック。「安全」「安楽」「自立支援」というラベル付き。青・緑・紫の3色グラデーション",
            "prompt_en": "Three large icons (shield, smiley face, handshake) arranged horizontally with labels Safety, Comfort, Independence Support. Blue, green, purple gradient, clean flat design, white background",
            "script": (
                "すべての介助に共通する根本的な考え方が「3大原則」です。\n\n"
                "1つ目は「安全」。転倒・転落・皮膚損傷を絶対に起こさない。\n"
                "「大丈夫だろう」は禁物です。確認してから動く。\n\n"
                "2つ目は「安楽」。利用者に苦痛や不安を与えない。\n"
                "急な動き・無言の介助はNGです。必ず声をかけてから動きましょう。\n\n"
                "3つ目は「自立支援」。全部やってしまわない。\n"
                "できることは本人にやってもらう。それが機能維持につながります。\n\n"
                "この3つを頭に入れて、全ての介助に臨んでください。"
            )
        },
        {
            "no": 4,
            "title": "ボディメカニクスとは",
            "time": "3分",
            "prompt_ja": "介護士が膝を曲げて低い重心で利用者を支えている場面。矢印や力の方向を示すグラフィックが重なっている。教育用のイラスト風、明るい背景",
            "prompt_en": "Caregiver with bent knees and low center of gravity supporting patient, arrows showing force direction and body mechanics principles, educational illustration style, bright background",
            "script": (
                "ボディメカニクスとは、人体の構造と物理の法則を活かして、\n"
                "介助者の体を守りながら安全に介助する技術です。\n\n"
                "腰痛は介護職の離職理由ナンバーワンです。\n"
                "正しい姿勢を身につけることが、長く働き続けるための第一歩です。\n\n"
                "原則は6つ。\n"
                "①足を肩幅に開く　②膝を曲げて重心を低くする\n"
                "③利用者に近づく　④大きな筋肉を使う\n"
                "⑤体をねじらない　⑥水平移動を活用する\n\n"
                "特に大事なのは「ねじり動作を避ける」こと。\n"
                "方向転換するときは、足ごと体を向けてください。"
            )
        },
        {
            "no": 5,
            "title": "移乗前チェックリスト",
            "time": "2分",
            "prompt_ja": "チェックリストのクリップボードと車椅子・ベッドのアイコンが並んだフラットイラスト。チェックマーク付き。青とオレンジのカラーアクセント",
            "prompt_en": "Flat illustration of checklist clipboard with wheelchair and bed icons, checkmarks, blue and orange color accents, clean modern design, infographic style",
            "script": (
                "移乗を始める前に必ず確認することがあります。\n\n"
                "車椅子のブレーキはかかっていますか？\n"
                "フットサポートは跳ね上げていますか？\n"
                "ベッドの高さは適切ですか？\n"
                "利用者への説明と同意は取れていますか？\n\n"
                "この確認を省くと事故につながります。\n"
                "「ブレーキよし！」と声に出して確認する習慣をつけてください。\n\n"
                "準備8割、介助2割——と覚えてください。"
            )
        },
        {
            "no": 6,
            "title": "ベッド→車椅子 手順①〜④",
            "time": "4分",
            "prompt_ja": "介護士が端座位をとった高齢者の横に立ち、前傾姿勢を促している場面。ステップ番号付きの吹き出しが画像に重なっている。明るい病室、フォトリアリスティック",
            "prompt_en": "Caregiver standing beside elderly person in sitting edge-of-bed position, guiding forward lean, step number callouts overlaying image, bright hospital room, photorealistic",
            "script": (
                "では実際の手順を確認しましょう。\n\n"
                "ステップ1：ベッドの高さを調整して端座位をとってもらいます。\n"
                "両足が床にしっかりつく高さが目安です。\n\n"
                "ステップ2：車椅子をベッドに対して20〜30度の角度で置きます。\n"
                "必ず健側（麻痺がない方）に車椅子を置いてください。\n\n"
                "ステップ3：「鼻がひざの上に来るように」前傾姿勢をとってもらいます。\n"
                "重心が前に移動することで立ちやすくなります。\n\n"
                "ステップ4：介助者は患側に立ち、腰に手を回します。\n"
                "膝で利用者の膝をブロックするように立つのがポイントです。"
            )
        },
        {
            "no": 7,
            "title": "ベッド→車椅子 手順⑤〜⑧",
            "time": "3分",
            "prompt_ja": "介護士と高齢者が一緒に立ち上がり、方向転換して車椅子に着座する連続動作を矢印で示した図解。4コマのシーケンス図、フラットイラスト",
            "prompt_en": "4-panel sequence diagram showing caregiver and elderly person standing up together, pivoting, and sitting in wheelchair with directional arrows. Flat illustration style",
            "script": (
                "ステップ5：「せーの」の声かけとともに一緒に立ち上がります。\n"
                "介助者は「膝を伸ばして持ち上げる」のではなく、\n"
                "利用者の動きに合わせてサポートするイメージです。\n\n"
                "ステップ6：立位が安定したら、健側方向に方向転換します。\n"
                "足をすり足で動かし、体をねじらないように。\n\n"
                "ステップ7：座面を確認してから、ゆっくり着座させます。\n"
                "深く座れているか必ず確認してください。\n\n"
                "ステップ8：フットサポートを戻し、「座れましたか？痛いところはないですか？」\n"
                "と声をかけて完了です。"
            )
        },
        {
            "no": 8,
            "title": "トイレ介助のポイント",
            "time": "3分",
            "prompt_ja": "バリアフリートイレの手すり・車椅子スペースを示した間取り図風のアイコンイラスト。プライバシーカーテンとナースコールボタンが強調されている",
            "prompt_en": "Floor plan icon illustration of accessible toilet showing grab bars, wheelchair space, privacy curtain and nurse call button highlighted, flat design, clean lines",
            "script": (
                "トイレ介助は利用者の尊厳に最も深く関わる介助です。\n\n"
                "まず入室前に床の水濡れを確認してください。\n"
                "次に車椅子を健側の手すり側に斜め45度で配置します。\n\n"
                "着座後は必ずカーテンを引き、\n"
                "「終わったらナースコールを押してください」と伝えてから退出します。\n\n"
                "大切なのは「プライバシーの確保」と「待たせないこと」。\n"
                "動ける準備を事前に整えておくことが重要です。\n\n"
                "失禁があっても、表情・声で驚いたり責めたりしないでください。\n"
                "利用者の心理的安全がケアの質を左右します。"
            )
        },
        {
            "no": 9,
            "title": "歩行介助の基本",
            "time": "3分",
            "prompt_ja": "介護士が高齢者の患側斜め後ろに立ち、腰に手を添えて廊下を歩いている場面。足の位置と手の添え方に注意書きが入っている教育用フォト",
            "prompt_en": "Caregiver standing diagonally behind elderly patient's affected side, hand placed on lower back while walking down corridor, educational annotations showing hand and foot placement, photorealistic",
            "script": (
                "歩行介助の目標は「転倒させない」だけではありません。\n"
                "「自分で歩ける自信をつけてもらう」ことも重要な目標です。\n\n"
                "立ち位置は患側の斜め後ろ。\n"
                "手は腰〜骨盤に添えます。脇を抱えるのは絶対NGです。\n"
                "肩関節脱臼のリスクがあります。\n\n"
                "利用者のペースに合わせて、焦らせないことが大切です。\n"
                "「段差があります」「もう少しで曲がります」と先に声かけしましょう。\n\n"
                "歩き方の変化や疲れの徴候も常に観察してください。"
            )
        },
        {
            "no": 10,
            "title": "車椅子移動・段差対応",
            "time": "3分",
            "prompt_ja": "介護士が車椅子を後ろ向きにして段差を降りている場面。赤い矢印で方向が示されている。安全な操作方法の教育用イラスト",
            "prompt_en": "Caregiver pushing wheelchair backward down a step/ramp with red directional arrows showing correct technique, safety instruction illustration, clean educational style",
            "script": (
                "車椅子で段差を下りるときは、必ず後ろ向きにしてください。\n"
                "前向きに降りると利用者が前傾して非常に危険です。\n\n"
                "スロープを上るときは前向きで、一定のスピードを保ちます。\n"
                "急停止はふらつきの原因になります。\n\n"
                "エレベーターは後ろ向きで乗り込み、前向きで降ります。\n\n"
                "走行中は常に声かけを忘れずに。\n"
                "利用者の目線では前方が見えにくいので、\n"
                "「曲がります」「段差があります」と先に伝えることが安心につながります。"
            )
        },
        {
            "no": 11,
            "title": "やってはいけないNG行為",
            "time": "2分",
            "prompt_ja": "バツマーク（❌）が重なった5つのNG介助シーンのアイコン図解。赤とグレーのカラースキーム、フラットイラスト",
            "prompt_en": "5 wrong caregiving posture icons with red X marks overlaid, showing incorrect techniques, red and gray color scheme, flat illustration style, educational warning design",
            "script": (
                "絶対にやってはいけないNG行為を確認しましょう。\n\n"
                "①脇を抱えて引き上げる → 肩関節脱臼のリスク\n"
                "②腕だけで引っ張る → 介助者の腰・肩を壊す\n"
                "③ブレーキをかけずに移乗 → 転倒事故の直接原因\n"
                "④前向きで段差を下りる → 前方転落の危険\n"
                "⑤急がせる・無言で動かす → 利用者の不安と転倒リスク増大\n\n"
                "一つひとつは「ちょっとしたミス」に見えますが、\n"
                "重大事故につながります。習慣化する前に今日から直しましょう。"
            )
        },
        {
            "no": 12,
            "title": "安全管理・ヒヤリハット",
            "time": "2分",
            "prompt_ja": "ヒヤリハット報告書のクリップボードと！マークの三角アイコン。チームで報告書を見ているスタッフの背景写真。ブルー＆オレンジのアクセントカラー",
            "prompt_en": "Clipboard with near-miss incident report form, warning triangle icon, background of staff reviewing documents together, blue and orange accent colors, photorealistic blend",
            "script": (
                "転倒や怪我は突然起こるわけではありません。\n"
                "必ずその前に「ヒヤリハット」があります。\n\n"
                "「大事にならなかったから報告しなくていい」は大きな間違いです。\n"
                "小さな気づきを報告・共有することで、次の大きな事故を防げます。\n\n"
                "施設のヒヤリハット報告書を積極的に活用してください。\n"
                "あなたの報告が、チーム全体を守ります。"
            )
        },
        {
            "no": 13,
            "title": "まとめ＆チェックリスト",
            "time": "2分",
            "prompt_ja": "若い看護師・介護士が笑顔でチェックリストを持っている場面。明るいオフィス/施設背景。達成感・前向きなムード",
            "prompt_en": "Young nurse and caregiver smiling while holding checklists, bright facility office background, sense of accomplishment and positive mood, photorealistic",
            "script": (
                "今日の研修のまとめです。\n\n"
                "・介助の3原則：安全・安楽・自立支援\n"
                "・ボディメカニクスで自分の体を守る\n"
                "・移乗前は必ずブレーキ確認\n"
                "・患側の斜め後ろに立って歩行介助\n"
                "・ヒヤリハットは必ず報告\n\n"
                "今日学んだことを明日からすぐに実践してください。\n"
                "わからないことや不安なことは先輩スタッフに必ず聞いてください。\n\n"
                "皆さんが安全で丁寧な介助ができるよう、応援しています。\n"
                "今日はお疲れ様でした！"
            )
        },
    ]

    for slide in slides:
        # スライドタイトル
        p = doc.add_paragraph()
        r = p.add_run(f"▶ スライド {slide['no']}：{slide['title']}　（目安 {slide['time']}）")
        r.bold = True; r.font.size = Pt(12)
        r.font.color.rgb = C_PRIMARY

        # 画像プロンプトボックス
        add_image_prompt_box(doc, slide['no'], slide['prompt_ja'], slide['prompt_en'])
        doc.add_paragraph()

        # 発表原稿
        p_script = doc.add_paragraph()
        r_s = p_script.add_run(slide['script'])
        r_s.font.size = Pt(10)

        doc.add_paragraph()
        # 区切り線
        p_hr = doc.add_paragraph()
        r_hr = p_hr.add_run("─" * 60)
        r_hr.font.size = Pt(8)
        r_hr.font.color.rgb = C_GRAY
        doc.add_paragraph()

    doc.save("C:/Users/yuuna/agens/新人勉強会/発表原稿_画像プロンプト付き.docx")
    print("OK: Script + image prompts saved")


# ═══════════════════════════════════════
#  3. PowerPoint スライド
# ═══════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm as PCm

PC_PRIMARY  = PRGBColor(0x1A, 0x56, 0xDB)
PC_ACCENT   = PRGBColor(0x0E, 0x9F, 0x6E)
PC_WARNING  = PRGBColor(0xE3, 0xA0, 0x08)
PC_DANGER   = PRGBColor(0xE0, 0x28, 0x28)
PC_WHITE    = PRGBColor(0xFF, 0xFF, 0xFF)
PC_DARK     = PRGBColor(0x1F, 0x29, 0x37)
PC_LIGHT    = PRGBColor(0xEB, 0xF5, 0xFF)
PC_GRAY     = PRGBColor(0x6B, 0x72, 0x80)
PC_BG       = PRGBColor(0xF8, 0xFA, 0xFF)

def add_shape(slide, left, top, width, height, fill_rgb=None, line_rgb=None, text="",
              font_size=18, font_bold=False, font_color=None, align=PP_ALIGN.LEFT, line_width=0):
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_width:
        shape.line.color.rgb = line_rgb
        shape.line.width = Emu(line_width)
    else:
        shape.line.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPt(font_size)
    run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = font_color
    return shape

def set_slide_bg(slide, r, g, b):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRGBColor(r, g, b)

def add_header_bar(slide, title, subtitle="", bg_color=PC_PRIMARY):
    W = Inches(10); H = Inches(9/16 * 1.5)
    # 背景バー
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = bg_color
    bar.line.fill.background()
    # タイトル
    tf = bar.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = f"  {title}"
    run.font.bold = True; run.font.size = PPt(28)
    run.font.color.rgb = PC_WHITE
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = f"  {subtitle}"
        r2.font.size = PPt(13); r2.font.color.rgb = PRGBColor(0xBF, 0xDB, 0xFF)

def add_bullet_box(slide, left, top, width, height, title, bullets,
                   title_color=PC_PRIMARY, bg=PC_LIGHT):
    # 枠
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.color.rgb = title_color; box.line.width = Emu(12700)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = PCm(0.3); tf.margin_top = PCm(0.2)
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run(); r0.text = title
    r0.font.bold = True; r0.font.size = PPt(14)
    r0.font.color.rgb = title_color
    for b in bullets:
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        p.space_before = PPt(2)
        r = p.add_run(); r.text = f"  ▸ {b}"
        r.font.size = PPt(11.5); r.font.color.rgb = PC_DARK

def add_step_box(slide, left, top, width, height, steps_list, col_colors=None):
    """番号付きステップを縦に並べる"""
    step_h = height // len(steps_list)
    for i, (num, txt) in enumerate(steps_list):
        y = top + i * step_h
        # 番号丸
        num_box = slide.shapes.add_shape(
            9,  # Oval
            left, y + Emu(50000),
            Emu(380000), Emu(380000)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = col_colors[i % len(col_colors)] if col_colors else PC_PRIMARY
        num_box.line.fill.background()
        num_tf = num_box.text_frame
        num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        nr = num_tf.paragraphs[0].add_run()
        nr.text = str(num)
        nr.font.size = PPt(12); nr.font.bold = True
        nr.font.color.rgb = PC_WHITE
        # テキスト
        txt_box = slide.shapes.add_textbox(
            left + Emu(450000), y + Emu(30000),
            width - Emu(500000), step_h - Emu(60000)
        )
        txt_box.fill.background(); txt_box.line.fill.background()
        tf = txt_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = txt
        r.font.size = PPt(11)


def create_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    blank_layout = prs.slide_layouts[6]  # blank

    # ── Slide 1: タイトル ──────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0x1A, 0x56, 0xDB)
    # 大タイトル
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "移乗・移動介助の基本"
    r.font.size = PPt(40); r.font.bold = True; r.font.color.rgb = PC_WHITE
    # サブ
    tb2 = sl.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.7))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Transfer & Mobility Assistance  ｜  新人研修"
    r2.font.size = PPt(16); r2.font.color.rgb = PRGBColor(0xBF, 0xDB, 0xFF)
    # 下部
    tb3 = sl.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "対象：看護・介護職員新入社員　｜　所要時間：約35分　｜　【画像をここに配置】"
    r3.font.size = PPt(11); r3.font.color.rgb = PRGBColor(0x93, 0xC5, 0xFD)
    # 画像プロンプトメモ
    img_note = sl.shapes.add_textbox(Inches(0.2), Inches(4.9), Inches(9.6), Inches(0.6))
    img_note.fill.solid(); img_note.fill.fore_color.rgb = PRGBColor(0x1E, 0x3A, 0x8A)
    img_note.line.fill.background()
    p_n = img_note.text_frame.paragraphs[0]; p_n.alignment = PP_ALIGN.CENTER
    r_n = p_n.add_run()
    r_n.text = "🎨 画像: 笑顔の看護師が車椅子移乗介助 / young nurse smiling assisting transfer, bright facility, warm light"
    r_n.font.size = PPt(8); r_n.font.color.rgb = PRGBColor(0x93, 0xC5, 0xFD)

    # ── Slide 2: 今日のゴール ─────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "今日のゴール", "この研修で身につける5つのこと")
    goals = [
        ("🛡️", "安全にできる", "ボディメカニクスで自分と利用者を守る"),
        ("🔄", "移乗できる", "ベッド↔車椅子の手順を確実に"),
        ("🚶", "歩行介助", "患側に寄り添い転倒を防ぐ"),
        ("⚡", "リスク予測", "転倒・転落を事前に察知する"),
        ("🤝", "自立支援", "できることは本人にやってもらう"),
    ]
    for i, (icon, title, desc) in enumerate(goals):
        x = Inches(0.3 + i * 1.92)
        card = sl.shapes.add_shape(1, x, Inches(1.6), Inches(1.75), Inches(3.5))
        card.fill.solid()
        colors = [
            PRGBColor(0xEB, 0xF5, 0xFF),
            PRGBColor(0xEC, 0xFD, 0xF5),
            PRGBColor(0xFE, 0xF3, 0xC7),
            PRGBColor(0xFD, 0xE8, 0xE8),
            PRGBColor(0xF3, 0xE8, 0xFF),
        ]
        card.fill.fore_color.rgb = colors[i]
        border_colors = [PC_PRIMARY, PC_ACCENT, PC_WARNING, PC_DANGER, PRGBColor(0x76, 0x1D, 0xA1)]
        card.line.color.rgb = border_colors[i]; card.line.width = Emu(25400)
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_top = PCm(0.3); tf.margin_left = PCm(0.2)
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = icon; r0.font.size = PPt(26)
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = title
        r1.font.bold = True; r1.font.size = PPt(12)
        r1.font.color.rgb = border_colors[i]
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = desc
        r2.font.size = PPt(9); r2.font.color.rgb = PC_GRAY

    img_note = sl.shapes.add_textbox(Inches(0.2), Inches(5.25), Inches(9.6), Inches(0.3))
    img_note.fill.solid(); img_note.fill.fore_color.rgb = PRGBColor(0xED, 0xEB, 0xFF)
    img_note.line.fill.background()
    p_n = img_note.text_frame.paragraphs[0]
    r_n = p_n.add_run()
    r_n.text = "🎨 5 flat icons: shield, arrows, walking figure, lightning bolt, handshake — blue/green/amber/red/purple"
    r_n.font.size = PPt(7.5); r_n.font.color.rgb = PRGBColor(0x4A, 0x1D, 0x96)

    # ── Slide 3: 介助の3大原則 ───────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "介助の3大原則", "すべての介助に共通する根本思想")
    for i, (title, desc, color) in enumerate([
        ("🛡️ 安全", "転倒・転落・皮膚損傷を防ぐ\nブレーキ確認・環境整備・声かけを徹底", PC_DANGER),
        ("😌 安楽", "利用者に苦痛・不安を与えない\n急な動き・無言の介助はNG", PC_ACCENT),
        ("🤝 自立支援", "できることは本人にやってもらう\n「やってあげる」ではなく「一緒にやる」", PC_PRIMARY),
    ]):
        x = Inches(0.4 + i * 3.1)
        box = sl.shapes.add_shape(1, x, Inches(1.7), Inches(2.9), Inches(3.5))
        box.fill.solid()
        bg_colors = [PRGBColor(0xFD, 0xE8, 0xE8), PRGBColor(0xEC, 0xFD, 0xF5), PRGBColor(0xEB, 0xF5, 0xFF)]
        box.fill.fore_color.rgb = bg_colors[i]
        box.line.color.rgb = color; box.line.width = Emu(38100)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_top = PCm(0.4); tf.margin_left = PCm(0.3)
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = title
        r0.font.bold = True; r0.font.size = PPt(18); r0.font.color.rgb = color
        p1 = tf.add_paragraph(); p1.space_before = PPt(8)
        r1 = p1.add_run(); r1.text = desc
        r1.font.size = PPt(11); r1.font.color.rgb = PC_DARK

    # ── Slide 4: ボディメカニクス ────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "ボディメカニクス 6原則", "介助者の体を守る科学的技術")
    bm = [
        "① 足を肩幅に開く（支持基底面を広げる）",
        "② 膝を曲げて重心を低くする",
        "③ 利用者にできるだけ近づく",
        "④ 大きな筋肉（下肢・体幹）を使う",
        "⑤ 体をねじらず足ごと方向転換",
        "⑥ 水平移動・てこの原理を活用する",
    ]
    col_c = [PC_PRIMARY, PC_ACCENT, PC_WARNING, PC_DANGER,
             PRGBColor(0x76, 0x1D, 0xA1), PRGBColor(0x03, 0x69, 0xA1)]
    for i, txt in enumerate(bm):
        row = i // 3; col = i % 3
        x = Inches(0.3 + col * 3.2)
        y = Inches(1.6 + row * 1.8)
        box = sl.shapes.add_shape(1, x, y, Inches(3.0), Inches(1.55))
        box.fill.solid(); box.fill.fore_color.rgb = PC_LIGHT
        box.line.color.rgb = col_c[i]; box.line.width = Emu(19050)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = PCm(0.3); tf.margin_top = PCm(0.2)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = txt
        r.font.size = PPt(12); r.font.bold = True; r.font.color.rgb = col_c[i]
    img_note = sl.shapes.add_textbox(Inches(9.5 - 3.0), Inches(1.6), Inches(3.0), Inches(3.5))
    # (image placeholder note in footer instead)

    # ── Slide 5: 移乗前チェック ──────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "移乗前チェック", "準備8割・介助2割！")
    checks = [
        "✅  車椅子のブレーキがかかっているか",
        "✅  フットサポートが跳ね上げられているか",
        "✅  ベッドの高さが適切か（足底が床に着く）",
        "✅  利用者に介助内容を説明し同意を得たか",
        "✅  周囲に障害物・危険物はないか",
        "✅  必要な補助具が手の届く場所にあるか",
    ]
    box = sl.shapes.add_shape(1, Inches(0.5), Inches(1.6), Inches(5.8), Inches(3.8))
    box.fill.solid(); box.fill.fore_color.rgb = PRGBColor(0xEC, 0xFD, 0xF5)
    box.line.color.rgb = PC_ACCENT; box.line.width = Emu(25400)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = PCm(0.4); tf.margin_top = PCm(0.3)
    for i, c in enumerate(checks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = c
        r.font.size = PPt(12); r.font.color.rgb = PC_DARK
    # 右側 警告ボックス
    warn_box = sl.shapes.add_shape(1, Inches(6.5), Inches(1.6), Inches(3.2), Inches(3.8))
    warn_box.fill.solid(); warn_box.fill.fore_color.rgb = PRGBColor(0xFF, 0xF3, 0xCD)
    warn_box.line.color.rgb = PC_WARNING; warn_box.line.width = Emu(25400)
    wf = warn_box.text_frame; wf.word_wrap = True
    wf.margin_left = PCm(0.3); wf.margin_top = PCm(0.3)
    p_w = wf.paragraphs[0]
    r_w = p_w.add_run(); r_w.text = "⚠️ 声に出して確認！"
    r_w.font.bold = True; r_w.font.size = PPt(14); r_w.font.color.rgb = PRGBColor(0x92, 0x40, 0x09)
    p_w2 = wf.add_paragraph(); p_w2.space_before = PPt(6)
    r_w2 = p_w2.add_run()
    r_w2.text = "「ブレーキよし！」\n「フットサポートよし！」\n\n声に出す習慣が\n事故を防ぎます"
    r_w2.font.size = PPt(11); r_w2.font.color.rgb = PC_DARK

    # ── Slide 6: ベッド→車椅子 手順 ──────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "ベッド → 車椅子　移乗手順", "8ステップで完全マスター")
    steps = [
        ("①", "ベッド高さ調整・端座位"),
        ("②", "車椅子を健側に20〜30°で配置"),
        ("③", "前傾姿勢（鼻がひざの上）"),
        ("④", "患側に立ち腰に手を添える"),
        ("⑤", "「せーの」で立ち上がり"),
        ("⑥", "健側へ方向転換"),
        ("⑦", "座面確認→ゆっくり着座"),
        ("⑧", "フットサポート・声かけ確認"),
    ]
    step_colors = [PC_PRIMARY, PC_PRIMARY, PC_ACCENT, PC_ACCENT,
                   PC_WARNING, PC_WARNING, PC_DANGER, PC_ACCENT]
    for i, (num, txt) in enumerate(steps):
        row = i // 4; col = i % 4
        x = Inches(0.25 + col * 2.4)
        y = Inches(1.55 + row * 1.85)
        box = sl.shapes.add_shape(1, x, y, Inches(2.25), Inches(1.7))
        box.fill.solid(); box.fill.fore_color.rgb = PC_LIGHT
        box.line.color.rgb = step_colors[i]; box.line.width = Emu(19050)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_top = PCm(0.15); tf.margin_left = PCm(0.2)
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = num
        r0.font.size = PPt(22); r0.font.bold = True; r0.font.color.rgb = step_colors[i]
        p1 = tf.add_paragraph()
        r1 = p1.add_run(); r1.text = txt
        r1.font.size = PPt(10.5); r1.font.color.rgb = PC_DARK

    # ── Slide 7: トイレ介助 ──────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "トイレ介助のポイント", "尊厳への配慮が最重要")
    add_bullet_box(sl, Inches(0.3), Inches(1.6), Inches(4.7), Inches(3.8),
                   "手順のポイント",
                   ["床の水濡れ確認（転倒防止）",
                    "健側手すり側・斜め45°配置",
                    "ブレーキ確認→フットサポート跳ね上げ",
                    "衣服は本人にやってもらう",
                    "カーテン引いてナースコール位置を伝える",
                    "終了後：手洗い介助・車椅子へ"],
                   title_color=PC_ACCENT)
    add_bullet_box(sl, Inches(5.2), Inches(1.6), Inches(4.5), Inches(3.8),
                   "尊厳への配慮",
                   ["「やってあげる」ではなく「一緒にやる」",
                    "無駄に滞在しない・プライバシー確保",
                    "待たせない（事前準備が重要）",
                    "失禁があっても驚かない・責めない"],
                   title_color=PRGBColor(0x76, 0x1D, 0xA1),
                   bg=PRGBColor(0xF3, 0xE8, 0xFF))

    # ── Slide 8: 歩行介助 ────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "歩行介助の基本", "転ばせない × 自信をつける")
    add_bullet_box(sl, Inches(0.3), Inches(1.6), Inches(4.7), Inches(3.8),
                   "立ち位置と手の添え方",
                   ["患側の斜め後ろに立つ",
                    "腰〜骨盤に手を添える",
                    "脇を抱えるのは絶対NG（脱臼リスク）",
                    "利用者のペースに合わせる",
                    "「段差あります」と先に声かけ"],
                   title_color=PC_ACCENT)
    # 補助具テーブル
    tb_box = sl.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.5), Inches(0.5))
    tf_b = tb_box.text_frame
    r_b = tf_b.paragraphs[0].add_run()
    r_b.text = "補助具使用時の段差対応"
    r_b.font.bold = True; r_b.font.size = PPt(13); r_b.font.color.rgb = PC_PRIMARY
    rows_data = [
        ("杖", "杖→健足→患足", "杖→患足→健足"),
        ("歩行器", "歩行器→両足", "歩行器→両足"),
    ]
    headers = ["補助具", "上り", "下り"]
    for r_i, row_data in enumerate([headers] + rows_data):
        for c_i, val in enumerate(row_data):
            bx = sl.shapes.add_shape(1,
                Inches(5.2 + c_i * 1.5), Inches(2.1 + r_i * 0.7),
                Inches(1.45), Inches(0.65))
            is_hdr = r_i == 0
            bx.fill.solid()
            bx.fill.fore_color.rgb = PC_PRIMARY if is_hdr else PC_LIGHT
            bx.line.color.rgb = PC_PRIMARY; bx.line.width = Emu(12700)
            tf_c = bx.text_frame; tf_c.word_wrap = True
            p_c = tf_c.paragraphs[0]; p_c.alignment = PP_ALIGN.CENTER
            r_c = p_c.add_run(); r_c.text = val
            r_c.font.size = PPt(9.5); r_c.font.bold = is_hdr
            r_c.font.color.rgb = PC_WHITE if is_hdr else PC_DARK

    # ── Slide 9: 車椅子移動・段差 ────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "車椅子移動介助", "段差・スロープの正しい操作")
    points = [
        ("🔽 段差を下りる", "必ず後ろ向き！\n前向きで降りると前傾・転落の危険", PC_DANGER),
        ("🔼 スロープを上る", "前向き・一定スピードで\n急停止はふらつきの原因", PC_ACCENT),
        ("🚪 エレベーター", "後ろ向きで乗り込み\n前向きで降りる", PC_PRIMARY),
        ("💬 走行中の声かけ", "「曲がります」「段差あります」\n常にアナウンスを", PRGBColor(0x76, 0x1D, 0xA1)),
    ]
    for i, (title, desc, color) in enumerate(points):
        x = Inches(0.3 + (i % 2) * 4.8)
        y = Inches(1.55 + (i // 2) * 1.95)
        box = sl.shapes.add_shape(1, x, y, Inches(4.4), Inches(1.75))
        box.fill.solid(); box.fill.fore_color.rgb = PC_LIGHT
        box.line.color.rgb = color; box.line.width = Emu(25400)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = PCm(0.3); tf.margin_top = PCm(0.2)
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = title
        r0.font.bold = True; r0.font.size = PPt(14); r0.font.color.rgb = color
        p1 = tf.add_paragraph(); p1.space_before = PPt(4)
        r1 = p1.add_run(); r1.text = desc
        r1.font.size = PPt(11); r1.font.color.rgb = PC_DARK

    # ── Slide 10: NG行為集 ───────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "絶対にやってはいけない NG行為", "", bg_color=PC_DANGER)
    ng_items = [
        ("❌ 脇を抱えて引き上げる", "→ 肩関節脱臼リスク"),
        ("❌ 腕だけで引っ張る", "→ 介助者の腰・肩を壊す"),
        ("❌ ブレーキをかけずに移乗", "→ 転倒事故の直接原因"),
        ("❌ 前向きで段差を下りる", "→ 前方転落の危険"),
        ("❌ 急がせる・無言で動かす", "→ 不安増大＋転倒リスク増"),
    ]
    for i, (ng, result) in enumerate(ng_items):
        row = i // 3; col = i % 3 if i < 3 else i - 3
        if i >= 3:
            x = Inches(1.7 + col * 3.3)
            y = Inches(3.2)
            w = Inches(3.1)
        else:
            x = Inches(0.3 + col * 3.2)
            y = Inches(1.55)
            w = Inches(3.0)
        box = sl.shapes.add_shape(1, x, y, w, Inches(1.5))
        box.fill.solid(); box.fill.fore_color.rgb = PRGBColor(0xFD, 0xE8, 0xE8)
        box.line.color.rgb = PC_DANGER; box.line.width = Emu(19050)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = PCm(0.3); tf.margin_top = PCm(0.2)
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = ng
        r0.font.bold = True; r0.font.size = PPt(11); r0.font.color.rgb = PC_DANGER
        p1 = tf.add_paragraph()
        r1 = p1.add_run(); r1.text = result
        r1.font.size = PPt(10); r1.font.color.rgb = PC_DARK

    # ── Slide 11: 安全管理・ヒヤリハット ─────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "安全管理とヒヤリハット報告", "小さな気づきが大きな事故を防ぐ")
    add_bullet_box(sl, Inches(0.3), Inches(1.6), Inches(4.7), Inches(3.8),
                   "主なリスクと予防策",
                   ["転倒：ブレーキ声出し確認・二人介助",
                    "腰痛：ボディメカニクス徹底・補助具",
                    "皮膚損傷：水平移動・スライディングシート",
                    "肩関節脱臼：骨盤・体幹を支える",
                    "誤嚥：食後30分は安静"],
                   title_color=PC_DANGER, bg=PRGBColor(0xFD, 0xE8, 0xE8))
    warn_box = sl.shapes.add_shape(1, Inches(5.3), Inches(1.6), Inches(4.3), Inches(3.8))
    warn_box.fill.solid(); warn_box.fill.fore_color.rgb = PRGBColor(0xFF, 0xF3, 0xCD)
    warn_box.line.color.rgb = PC_WARNING; warn_box.line.width = Emu(38100)
    wf = warn_box.text_frame; wf.word_wrap = True
    wf.margin_left = PCm(0.4); wf.margin_top = PCm(0.3)
    for txt in [
        "⚠️ ヒヤリハット報告",
        "",
        "「大事にならなかったから\n報告しなくていい」は❌",
        "",
        "小さな気づきを報告・共有が\n次の大事故を防ぐ",
        "",
        "あなたの報告が\nチーム全体を守る！",
    ]:
        p = wf.paragraphs[0] if txt == "⚠️ ヒヤリハット報告" else wf.add_paragraph()
        r = p.add_run(); r.text = txt
        if txt == "⚠️ ヒヤリハット報告":
            r.font.bold = True; r.font.size = PPt(15); r.font.color.rgb = PRGBColor(0x92, 0x40, 0x09)
        else:
            r.font.size = PPt(11); r.font.color.rgb = PC_DARK

    # ── Slide 12: チェックリスト ─────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0xF8, 0xFA, 0xFF)
    add_header_bar(sl, "研修後 自己チェックリスト", "「わかった」を「できる」に変えよう")
    cl_items = [
        "□ ボディメカニクス6原則を言える",
        "□ ベッド→車椅子の手順を8ステップで説明できる",
        "□ ブレーキ確認を声に出す習慣がついた",
        "□ 患側に立って歩行介助できる",
        "□ 段差での車椅子後ろ向き操作を理解した",
        "□ トイレ介助で尊厳への配慮ができる",
        "□ ヒヤリハットを必ず報告するという意識がある",
    ]
    box = sl.shapes.add_shape(1, Inches(0.5), Inches(1.55), Inches(9), Inches(3.7))
    box.fill.solid(); box.fill.fore_color.rgb = PRGBColor(0xEC, 0xFD, 0xF5)
    box.line.color.rgb = PC_ACCENT; box.line.width = Emu(25400)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = PCm(0.5); tf.margin_top = PCm(0.3)
    for i, item in enumerate(cl_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = item
        r.font.size = PPt(13); r.font.color.rgb = PC_DARK

    # ── Slide 13: まとめ ─────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, 0x1A, 0x56, 0xDB)
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "まとめ　—　今日学んだこと"
    r.font.size = PPt(26); r.font.bold = True; r.font.color.rgb = PC_WHITE

    summary_items = [
        ("🛡️", "3大原則", "安全・安楽・自立支援", PC_DANGER),
        ("💪", "ボディメカニクス", "6原則で腰痛ゼロへ", PC_ACCENT),
        ("🔄", "移乗手順", "ブレーキ確認が命綱", PC_WARNING),
        ("🚶", "歩行介助", "患側の斜め後ろに立つ", PRGBColor(0x06, 0xB6, 0xD4)),
        ("⚠️", "ヒヤリハット", "必ず報告・チームで共有", PRGBColor(0xA8, 0x55, 0xFF)),
    ]
    for i, (icon, title, desc, color) in enumerate(summary_items):
        col = i % 3 if i < 3 else i - 3 + (1 if i < 5 else 0)
        row = 0 if i < 3 else 1
        x = Inches(0.3 + (col if i < 3 else col + 0.5) * 3.1)
        y = Inches(1.3 + row * 2.0)
        box = sl.shapes.add_shape(1, x, y, Inches(2.85), Inches(1.75))
        box.fill.solid(); box.fill.fore_color.rgb = PRGBColor(0x1E, 0x3A, 0x8A)
        box.line.color.rgb = color; box.line.width = Emu(25400)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_top = PCm(0.25); tf.margin_left = PCm(0.3)
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = f"{icon}  {title}"
        r0.font.bold = True; r0.font.size = PPt(14); r0.font.color.rgb = color
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = desc
        r1.font.size = PPt(11); r1.font.color.rgb = PRGBColor(0xBF, 0xDB, 0xFF)

    # フッター
    ft = sl.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9), Inches(0.4))
    tf_f = ft.text_frame
    p_f = tf_f.paragraphs[0]; p_f.alignment = PP_ALIGN.CENTER
    r_f = p_f.add_run(); r_f.text = "明日から実践しましょう！わからないことは先輩スタッフに遠慮なく聞いてください。"
    r_f.font.size = PPt(11); r_f.font.color.rgb = PRGBColor(0x93, 0xC5, 0xFD)

    prs.save("C:/Users/yuuna/agens/新人勉強会/スライド_移乗移動介助.pptx")
    print("OK: PowerPoint saved")


if __name__ == "__main__":
    create_handout()
    create_script()
    create_pptx()
    print("DONE: All files created in C:/Users/yuuna/agens/新人勉強会/")
