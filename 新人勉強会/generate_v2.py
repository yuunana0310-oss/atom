# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Pt as DPt, Cm as DCm, RGBColor as DRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# ── PPTX colors ──────────────────────────────────────────
PC_NAVY   = RGBColor(0x1E,0x3A,0x8A)
PC_BLUE   = RGBColor(0x1D,0x4E,0xD8)
PC_TEAL   = RGBColor(0x06,0x82,0x9C)
PC_GREEN  = RGBColor(0x05,0x96,0x69)
PC_AMBER  = RGBColor(0xD9,0x77,0x06)
PC_RED    = RGBColor(0xDC,0x26,0x26)
PC_PURPLE = RGBColor(0x76,0x1D,0xA1)
PC_WHITE  = RGBColor(0xFF,0xFF,0xFF)
PC_DARK   = RGBColor(0x1F,0x29,0x37)
PC_GRAY   = RGBColor(0x6B,0x72,0x80)
PC_LIGHT  = RGBColor(0xEB,0xF5,0xFF)
PC_BG     = RGBColor(0xF8,0xFA,0xFF)
PC_LBLUE  = RGBColor(0x93,0xC5,0xFD)
PC_WARN_BG= RGBColor(0xFF,0xF3,0xCD)
PC_RED_BG = RGBColor(0xFD,0xE8,0xE8)
PC_GRN_BG = RGBColor(0xEC,0xFD,0xF5)
PC_PUR_BG = RGBColor(0xF3,0xE8,0xFF)
PC_LNAV   = RGBColor(0xBF,0xDB,0xFF)

W = Inches(13.33)
H = Inches(7.5)

# ── helpers ──────────────────────────────────────────────
def rect(slide, l,t,w,h, fill=None, line=None, lw=19050):
    s = slide.shapes.add_shape(1,l,t,w,h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*fill) if isinstance(fill,tuple) else fill
    else:    s.fill.background()
    if line: s.line.color.rgb = RGBColor(*line) if isinstance(line,tuple) else line; s.line.width=Emu(lw)
    else:    s.line.fill.background()
    return s

def oval(slide, l,t,w,h, fill=None):
    s = slide.shapes.add_shape(9,l,t,w,h)  # 9=oval
    if fill: s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*fill) if isinstance(fill,tuple) else fill
    else:    s.fill.background()
    s.line.fill.background()
    return s

def txbox(slide, l,t,w,h, text="", size=12, bold=False, color=PC_DARK,
          align=PP_ALIGN.LEFT, fill=None, line=None, lw=19050, wrap=True,
          italic=False, spacing_before=0):
    tb = slide.shapes.add_textbox(l,t,w,h)
    tf = tb.text_frame; tf.word_wrap=wrap
    if fill: tb.fill.solid(); tb.fill.fore_color.rgb=RGBColor(*fill) if isinstance(fill,tuple) else fill
    else:    tb.fill.background()
    if line: tb.line.color.rgb=RGBColor(*line) if isinstance(line,tuple) else line; tb.line.width=Emu(lw)
    else:    tb.line.fill.background()
    p = tf.paragraphs[0]; p.alignment=align
    if spacing_before: p.space_before=Pt(spacing_before)
    r = p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb = RGBColor(*color) if isinstance(color,tuple) else color
    return tb,tf

def add_para(tf, text, size=11, bold=False, color=PC_DARK, align=PP_ALIGN.LEFT, before=2):
    p = tf.add_paragraph(); p.alignment=align; p.space_before=Pt(before)
    r = p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb = RGBColor(*color) if isinstance(color,tuple) else color

def set_bg(slide, color):
    bg=slide.background; bg.fill.solid()
    bg.fill.fore_color.rgb=RGBColor(*color) if isinstance(color,tuple) else color

def header(slide, title, sub="", num="", bg=PC_NAVY):
    set_bg(slide, PC_BG)
    rect(slide, 0,0, W, Inches(1.4), fill=bg)
    txbox(slide, Inches(0.45),Inches(0.06),Inches(12.0),Inches(0.75),
          title,26,True,PC_WHITE)
    if sub:
        txbox(slide,Inches(0.45),Inches(0.78),Inches(11.5),Inches(0.45),
              sub,12,color=PC_LNAV)
    if num:
        txbox(slide,Inches(12.3),Inches(0.06),Inches(0.9),Inches(0.4),
              num,11,color=PC_LBLUE,align=PP_ALIGN.RIGHT)

def imgbox(slide, l,t,w,h, label="画像をここに配置"):
    rect(slide,l,t,w,h,fill=(0xF1,0xF5,0xF9),line=(0x94,0xA3,0xB8))
    txbox(slide,l+Inches(0.1),t+h//2-Inches(0.3),w-Inches(0.2),Inches(0.6),
          f"[ {label} ]",10,color=(0x94,0xA3,0xB8),align=PP_ALIGN.CENTER)

def infobox(slide,l,t,w,h,title,lines,tc=PC_NAVY,bg=PC_LIGHT,bc=None):
    if bc is None: bc=tc
    rect(slide,l,t,w,h,fill=bg,line=bc,lw=25400)
    txbox(slide,l+Inches(0.15),t+Inches(0.12),w-Inches(0.25),Inches(0.45),
          title,13,True,tc)
    tb,tf=txbox(slide,l+Inches(0.15),t+Inches(0.58),w-Inches(0.25),h-Inches(0.65),"",11,color=PC_DARK)
    tf.word_wrap=True
    for i,ln in enumerate(lines):
        if i==0: p=tf.paragraphs[0]
        else:    p=tf.add_paragraph(); p.space_before=Pt(3)
        r=p.add_run(); r.text=ln; r.font.size=Pt(11); r.font.color.rgb=PC_DARK

def warnbox(slide,l,t,w,h,text,danger=False):
    c = PC_RED if danger else PC_AMBER
    bg= PC_RED_BG if danger else PC_WARN_BG
    rect(slide,l,t,w,h,fill=bg,line=c,lw=25400)
    icon="!" if danger else "▲"
    txbox(slide,l+Inches(0.15),t+Inches(0.1),w-Inches(0.25),h-Inches(0.15),
          f"{icon}  {text}",11,color=PC_DARK)

def step_num(slide,l,t,sz,n,color=PC_NAVY):
    s=oval(slide,l,t,sz,sz,fill=color)
    tf=s.text_frame; tf.margin_top=Emu(0); tf.margin_left=Emu(0)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(n); r.font.size=Pt(12); r.font.bold=True
    r.font.color.rgb=PC_WHITE

# ════════════════════════════════════════════════════════
#  PPTX
# ════════════════════════════════════════════════════════
def make_pptx():
    prs = Presentation()
    prs.slide_width=W; prs.slide_height=H
    blank=prs.slide_layouts[6]

    # ── Slide 1: タイトル ─────────────────────────────
    sl=prs.slides.add_slide(blank); set_bg(sl,PC_NAVY)
    rect(sl,0,0,W,H,fill=PC_NAVY)
    # 斜めアクセント
    rect(sl,Inches(8.5),0,Inches(4.83),H,fill=(0x1D,0x4E,0xD8))
    txbox(sl,Inches(0.8),Inches(1.5),Inches(8.0),Inches(0.7),
          "新人職員研修",14,color=PC_LNAV)
    txbox(sl,Inches(0.8),Inches(2.1),Inches(8.2),Inches(1.6),
          "移乗・移動介助の基本",38,True,PC_WHITE)
    txbox(sl,Inches(0.8),Inches(3.75),Inches(8.0),Inches(0.6),
          "〜安全・安楽・自立支援を実践する〜",14,color=PC_LNAV)
    txbox(sl,Inches(0.8),Inches(4.7),Inches(8.0),Inches(0.55),
          "対象：看護職員・介護職員　｜　所要時間：約30分",12,color=PC_LBLUE)
    txbox(sl,Inches(0.8),Inches(5.4),Inches(8.0),Inches(0.45),
          "医療法人 研修・教育委員会",11,color=PC_LBLUE)
    imgbox(sl,Inches(9.0),Inches(1.2),Inches(3.8),Inches(4.8),
           "移乗介助の場面写真\nyoung nurse assisting transfer\nbright facility warm light")

    # ── Slide 2: 研修目標 ─────────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"本日の研修目標","この研修で身につけること","2 / 18")
    goals=[
        ("①",PC_BLUE,"移乗・移動介助の基本原則とボディメカニクスを理解する"),
        ("②",PC_TEAL,"車椅子・ベッド間の移乗手順を正しく実施できる"),
        ("③",PC_GREEN,"歩行介助・立ち上がり介助を安全に行える"),
        ("④",PC_AMBER,"転倒・転落リスクを予測し予防策を講じられる"),
        ("⑤",PC_PURPLE,"患者の尊厳と残存能力を尊重した介助ができる"),
    ]
    for i,(num,col,txt) in enumerate(goals):
        y=Inches(1.55)+i*Inches(1.1)
        h=str(col); cr,cg,cb=int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
        rect(sl,Inches(0.5),y,Inches(12.33),Inches(0.95),
             fill=(min(cr+40,255),min(cg+40,255),min(cb+40,255)),
             line=col,lw=19050)
        # 番号丸
        s=oval(sl,Inches(0.65),y+Inches(0.18),Inches(0.58),Inches(0.58),fill=col)
        tf=s.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=num; r.font.size=Pt(14); r.font.bold=True
        r.font.color.rgb=PC_WHITE
        txbox(sl,Inches(1.35),y+Inches(0.2),Inches(11.0),Inches(0.55),
              txt,14,color=PC_DARK)
    warnbox(sl,Inches(0.5),Inches(7.0),Inches(12.33),Inches(0.35),
            "患者の安全・安楽・自立支援が介助の3原則です")

    # ── Slide 3: 介助の3原則 ──────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"介助の3原則","すべての介助の土台となる考え方","3 / 18")
    cols3=[(PC_RED,(0xFD,0xE8,0xE8)),(PC_TEAL,PC_GRN_BG),(PC_PURPLE,PC_PUR_BG)]
    titles3=["安　全","安　楽","自立支援"]
    checks3=[
        ["転倒・転落・皮膚損傷を防ぐ","環境整備を必ず行う","2人介助の基準を遵守する"],
        ["患者の苦痛・不安を最小化","声かけ・同意を必ず得る","プライバシーを守る"],
        ["できることは自分で行ってもらう","残存能力を活かす介助","ADL維持・向上を目指す"],
    ]
    for i,(c,bg) in enumerate(cols3):
        x=Inches(0.4+i*4.3)
        rect(sl,x,Inches(1.55),Inches(4.1),Inches(5.3),fill=bg,line=c,lw=38100)
        txbox(sl,x+Inches(0.15),Inches(1.65),Inches(3.8),Inches(0.65),
              titles3[i],22,True,c,align=PP_ALIGN.CENTER)
        for j,ck in enumerate(checks3[i]):
            txbox(sl,x+Inches(0.2),Inches(2.4)+j*Inches(0.85),Inches(3.7),Inches(0.75),
                  f"✓  {ck}",12,color=PC_DARK)
    warnbox(sl,Inches(0.4),Inches(7.0),Inches(12.53),Inches(0.38),
            "声かけ・同意・プライバシー保護は必ず実施！")

    # ── Slide 4: ボディメカニクス ─────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"ボディメカニクスの基本","介助者の腰痛予防と安全な介助のために","4 / 18")
    bm=[
        (PC_BLUE, "① 支持基底面を広げる","足を肩幅に開き、安定した姿勢を確保する"),
        (PC_TEAL, "② 重心を低くする","膝を曲げて腰を落とし、重心を下げる"),
        (PC_GREEN,"③ 重心を近づける","患者との距離を縮め、密着して介助する"),
        (PC_AMBER,"④ 大きな筋群を使う","腰だけでなく脚・体幹の大きな筋肉を使う"),
        (PC_RED,  "⑤ 身体をひとつにまとめる","患者の手足を体に引き付けてから動かす"),
        (PC_PURPLE,"⑥ ねじり動作を避ける","足を向けてから体を回旋させる"),
    ]
    for i,(c,t,d) in enumerate(bm):
        col=i%2; row=i//2
        x=Inches(0.35+col*4.1); y=Inches(1.55)+row*Inches(1.9)
        rect(sl,x,y,Inches(3.9),Inches(1.75),fill=PC_BG,line=c,lw=25400)
        s=oval(sl,x+Inches(0.1),y+Inches(0.12),Inches(0.55),Inches(0.55),fill=c)
        p=s.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=str(i+1); r.font.size=Pt(12); r.font.bold=True
        r.font.color.rgb=PC_WHITE
        txbox(sl,x+Inches(0.75),y+Inches(0.12),Inches(3.0),Inches(0.5),t,12,True,c)
        txbox(sl,x+Inches(0.75),y+Inches(0.62),Inches(3.0),Inches(0.95),d,10,color=PC_GRAY)
    imgbox(sl,Inches(8.5),Inches(1.55),Inches(4.5),Inches(5.7),
           "【画像①】ボディメカニクス6原則\nイラスト（姿勢比較図）")
    warnbox(sl,Inches(0.35),Inches(7.0),Inches(12.53),Inches(0.38),
            "腰痛は職業病ではありません。正しい姿勢で自分の体も守りましょう")

    # ── Slide 5: 介助前確認 ───────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"介助前の確認事項","リスクアセスメントと環境整備","5 / 18")
    infobox(sl,Inches(0.35),Inches(1.55),Inches(6.0),Inches(5.1),
            "患者確認",
            ["● 患者の理解・同意を得る",
             "● バイタルサイン・体調確認",
             "● 点滴・ドレーン・カテーテル確認",
             "● 麻痺・疼痛・可動域の確認",
             "● 患者の体重・介助レベル確認"],
            tc=PC_TEAL,bg=PC_GRN_BG,bc=PC_TEAL)
    infobox(sl,Inches(6.55),Inches(1.55),Inches(6.44),Inches(5.1),
            "環境整備",
            ["● ベッドの高さ調整（介助しやすい高さ）",
             "● ブレーキ・ストッパーの確認",
             "● 床の濡れ・障害物の除去",
             "● スリッパ→滑り止め靴に変更",
             "● プライバシーカーテンを閉める"],
            tc=PC_NAVY,bg=PC_LIGHT,bc=PC_NAVY)
    warnbox(sl,Inches(0.35),Inches(6.75),Inches(12.63),Inches(0.62),
            "2人介助基準：全介助・体重60kg以上・不安定な患者は必ず2名で対応",danger=True)

    # ── Slide 6: 立ち上がり介助 ───────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"立ち上がり介助","ベッド端座位から立位へ","6 / 18")
    steps6=[
        (PC_BLUE, "STEP 1","端座位を確保する","ベッド端に座り、足底を床につける\n（足が浮く場合はフットレストを使用）"),
        (PC_TEAL, "STEP 2","前傾姿勢をとる","体を前に傾け重心を前方に移動\n介助者は腰部を支持する"),
        (PC_AMBER,"STEP 3","立ち上がりを誘導","「せーの」の声かけで一緒に立ち上がる\n脇・腰を支えながら立位へ"),
        (PC_GREEN,"STEP 4","立位の安定確認","めまい・ふらつきを確認\n必要なら手すり・歩行器を使用"),
    ]
    for i,(c,st,ti,de) in enumerate(steps6):
        x=Inches(0.35+i*3.2); y=Inches(1.55)
        rect(sl,x,y,Inches(3.0),Inches(5.2),fill=PC_BG,line=c,lw=25400)
        s=oval(sl,x+Inches(1.1),y+Inches(0.2),Inches(0.8),Inches(0.8),fill=c)
        p=s.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=st; r.font.size=Pt(10); r.font.bold=True
        r.font.color.rgb=PC_WHITE
        txbox(sl,x+Inches(0.1),y+Inches(1.2),Inches(2.8),Inches(0.55),ti,14,True,c,align=PP_ALIGN.CENTER)
        txbox(sl,x+Inches(0.15),y+Inches(1.85),Inches(2.7),Inches(3.1),de,11,color=PC_DARK)
    warnbox(sl,Inches(0.35),Inches(6.88),Inches(12.63),Inches(0.5),
            "ポイント：患者に合わせた速度で！急かすと転倒リスクが急増します")

    # ── Slide 7: ベッド→車椅子 ───────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"ベッド → 車椅子 移乗手順","片麻痺患者（健側を活用した移乗）","7 / 18")
    steps7=[
        "車椅子をベッドに対して約30°の角度で設置（健側に置く）",
        "車椅子のブレーキをかけ、フットレストを上げる",
        "患者を端座位にし、足底を床につける",
        "患者に前傾姿勢をとってもらい、健側の手で介助",
        "介助者は患側腰部〜臀部を支え、一緒に立ち上がる",
        "ゆっくり方向転換し、車椅子に向かって後ろに下がる",
        "臀部をシートに誘導しながらゆっくり着座させる",
        "フットレストを下げ、足を乗せて姿勢を整える",
    ]
    colors7=[PC_BLUE,PC_BLUE,PC_TEAL,PC_TEAL,PC_AMBER,PC_AMBER,PC_GREEN,PC_GREEN]
    for i,(txt,c) in enumerate(zip(steps7,colors7)):
        col=i%4; row=i//4
        x=Inches(0.35+col*2.35); y=Inches(1.55+row*2.6)
        rect(sl,x,y,Inches(2.2),Inches(2.4),fill=PC_BG,line=c,lw=19050)
        s=oval(sl,x+Inches(0.75),y+Inches(0.1),Inches(0.65),Inches(0.65),fill=c)
        p=s.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=str(i+1); r.font.size=Pt(14); r.font.bold=True
        r.font.color.rgb=PC_WHITE
        txbox(sl,x+Inches(0.1),y+Inches(0.85),Inches(2.0),Inches(1.45),txt,10,color=PC_DARK)
    imgbox(sl,Inches(9.5),Inches(1.55),Inches(3.5),Inches(2.45),
           "【画像②】ベッド→車椅子\n手順イラスト(STEP1〜4)")
    imgbox(sl,Inches(9.5),Inches(4.15),Inches(3.5),Inches(2.45),
           "【画像③】車椅子の配置\n角度図（30°配置）")

    # ── Slide 8: 車椅子→ベッド ───────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"車椅子 → ベッド 移乗手順","基本の流れと注意点","8 / 18")
    txbox(sl,Inches(0.5),Inches(1.5),Inches(12.33),Inches(0.4),
          "車椅子→ベッドは逆の手順で行います。ポイントを確認しましょう。",
          13,color=PC_GRAY)
    items8=[
        (PC_NAVY, "車椅子の配置",
         "ベッドに対して健側に30°で設置。\nブレーキ・フットレスト確認を先に行う"),
        (PC_TEAL, "体重移動の誘導",
         "「前に体を倒して」と声かけ。\n前傾姿勢から立ち上がりを誘導"),
        (PC_AMBER,"方向転換",
         "介助者は患者の患側に位置。\n転倒防止のため体を密着させる"),
        (PC_GREEN,"着座誘導",
         "ベッドに手をつかせてから、\nゆっくり膝を曲げて座らせる"),
        (PC_PURPLE,"体位の整え",
         "座位バランス確認後、\nベッド上での体位を整える"),
    ]
    bgs8=[(0xEB,0xF5,0xFF),(0xEC,0xFD,0xF5),(0xFF,0xF3,0xCD),(0xEC,0xFD,0xF5),(0xF3,0xE8,0xFF)]
    for i,((c,t,d),bg) in enumerate(zip(items8,bgs8)):
        x=Inches(0.35+i*2.5); y=Inches(1.95)
        rect(sl,x,y,Inches(2.35),Inches(4.8),fill=bg,line=c,lw=25400)
        rect(sl,x,y,Inches(2.35),Inches(0.55),fill=c)
        txbox(sl,x+Inches(0.1),y+Inches(0.06),Inches(2.15),Inches(0.45),t,12,True,PC_WHITE)
        txbox(sl,x+Inches(0.12),y+Inches(0.7),Inches(2.1),Inches(3.9),d,11,color=PC_DARK)
    warnbox(sl,Inches(0.35),Inches(6.88),Inches(12.63),Inches(0.5),
            "両移乗共通：患者の「できること」を最大限活かす介助を心がけましょう")

    # ── Slide 9: 移乗補助用具 ────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"移乗補助用具の活用","スライディングボード・移乗ベルト・リフト","9 / 18")
    tools9=[
        (PC_BLUE, "スライディングボード",
         ["・座位移乗に使用","・車椅子⇔ベッド・シャワーチェア",
          "・麻痺・筋力低下患者に有効","・使用時は必ず2名で確認"],
         "【画像④】スライディングボードの使用イラスト"),
        (PC_TEAL, "移乗ベルト（トランスファーベルト）",
         ["・患者の腰に装着","・介助者がベルトを握って誘導",
          "・急な体重移動に対応しやすい","・患者の転倒防止に有効"],
         "【画像⑤】移乗ベルト装着・使用方法イラスト"),
        (PC_PURPLE,"介護リフト",
         ["・全介助・重度障害患者に使用","・スリングシートを正しく装着",
          "・事前にスリング選択が重要","・操作は必ず研修後に実施"],
         "【画像⑥】介護リフト操作イラスト"),
    ]
    for i,(c,t,ls,img) in enumerate(tools9):
        x=Inches(0.35+i*4.3)
        imgbox(sl,x,Inches(1.55),Inches(4.1),Inches(2.8),img)
        rect(sl,x,Inches(4.45),Inches(4.1),Inches(2.75),fill=PC_BG,line=c,lw=19050)
        txbox(sl,x+Inches(0.1),Inches(4.5),Inches(3.9),Inches(0.45),t,12,True,c)
        tb,tf=txbox(sl,x+Inches(0.1),Inches(5.0),Inches(3.9),Inches(2.0),"",11,color=PC_DARK)
        tf.word_wrap=True
        for j,ln in enumerate(ls):
            p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
            r=p.add_run(); r.text=ln; r.font.size=Pt(11); r.font.color.rgb=PC_DARK
    warnbox(sl,Inches(0.35),Inches(7.2),Inches(12.63),Inches(0.25),
            "用具の使用前は必ず破損・劣化チェックを行うこと")

    # ── Slide 10: 歩行介助 ───────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"歩行介助の基本","安全な歩行介助の位置と手順","10 / 18")
    txbox(sl,Inches(0.35),Inches(1.5),Inches(12.63),Inches(0.4),
          "介助者の位置：患側（麻痺のある側）の斜め後方 ／ 片手：患者の腰または移乗ベルト ／ 歩幅・速度は患者に合わせる",
          11,color=PC_GRAY)
    phases10=[
        (PC_NAVY, "開始前",  ["靴・履物確認","バイタル確認","歩行補助具の用意"]),
        (PC_BLUE, "出発時",  ["ベッドサイドで立位保持","しばらく安静","めまい確認"]),
        (PC_TEAL, "歩行中",  ["声かけしながら","一定ペースで","障害物に注意"]),
        (PC_AMBER,"方向転換",["健側に向かって","ゆっくり回転","急な向き変え禁止"]),
        (PC_GREEN,"着座時",  ["椅子/ベッドに触れてから","ゆっくり膝を曲げる","確実に着座"]),
    ]
    for i,(c,t,ls) in enumerate(phases10):
        x=Inches(0.35+i*2.42)
        rect(sl,x,Inches(1.98),Inches(2.28),Inches(3.8),fill=PC_BG,line=c,lw=25400)
        rect(sl,x,Inches(1.98),Inches(2.28),Inches(0.55),fill=c)
        txbox(sl,x+Inches(0.07),Inches(2.02),Inches(2.14),Inches(0.45),t,13,True,PC_WHITE,align=PP_ALIGN.CENTER)
        tb,tf=txbox(sl,x+Inches(0.1),Inches(2.62),Inches(2.08),Inches(3.0),"",11,color=PC_DARK)
        tf.word_wrap=True
        for j,ln in enumerate(ls):
            p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
            p.space_before=Pt(4)
            r=p.add_run(); r.text=f"• {ln}"; r.font.size=Pt(11); r.font.color.rgb=PC_DARK
    imgbox(sl,Inches(0.5),Inches(5.88),Inches(5.8),Inches(1.35),
           "【画像⑦】歩行介助の位置・手の添え方イラスト")
    imgbox(sl,Inches(6.6),Inches(5.88),Inches(6.38),Inches(1.35),
           "【画像⑧】歩行補助具（T字杖・歩行器・ロフストランドクラッチ）の写真")

    # ── Slide 11: 歩行補助具 ─────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"歩行補助具の種類と特徴","患者の状態に合った補助具を選択する","11 / 18")
    aids11=[
        (PC_BLUE, "T字杖（ステッキ）",
         ["軽度の歩行障害に適応","健側で使用する","4点杖より安定性は低い"],
         "【写真】T字杖"),
        (PC_TEAL, "4点杖",
         ["接地面が広く安定","片麻痺・バランス障害に有効","室内使用が主"],
         "【写真】4点杖"),
        (PC_AMBER,"歩行器（固定型）",
         ["前方支持で安定","両上肢に体重をかけられる","段差・屋外には不向き"],
         "【写真】歩行器"),
        (PC_GREEN,"ロールウォーカー",
         ["車輪付きで歩きやすい","屋内外兼用可","ブレーキ操作を確認"],
         "【写真】ロールウォーカー"),
    ]
    for i,(c,t,ls,img) in enumerate(aids11):
        x=Inches(0.35+i*3.2)
        imgbox(sl,x,Inches(1.55),Inches(3.0),Inches(2.5),img)
        rect(sl,x,Inches(4.1),Inches(3.0),Inches(3.1),fill=PC_BG,line=c,lw=19050)
        txbox(sl,x+Inches(0.1),Inches(4.15),Inches(2.8),Inches(0.5),t,13,True,c)
        tb,tf=txbox(sl,x+Inches(0.1),Inches(4.72),Inches(2.8),Inches(2.3),"",11,color=PC_DARK)
        tf.word_wrap=True
        for j,ln in enumerate(ls):
            p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
            r=p.add_run(); r.text=f"・{ln}"; r.font.size=Pt(11); r.font.color.rgb=PC_DARK
    warnbox(sl,Inches(0.35),Inches(7.2),Inches(12.63),Inches(0.27),
            "補助具の選択は理学療法士・主治医と連携して決定すること")

    # ── Slide 12: 全介助（2名介助） ──────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"全介助の移乗（2名介助）","意識障害・全身脱力・重度障害患者への対応","12 / 18")
    warnbox(sl,Inches(0.35),Inches(1.5),Inches(12.63),Inches(0.55),
            "2名介助の適応：全介助患者 ／ 体重60kg以上 ／ 不安定な患者 ／ 初回移乗",danger=True)
    infobox(sl,Inches(0.35),Inches(2.15),Inches(5.9),Inches(4.95),
            "リーダー役（頭側）",
            ["✓  患者の上半身・頭頸部を支持",
             "✓  声かけ・タイミングの指示を出す",
             "✓  「いちにのさん」で全員同時に動く",
             "✓  気道・カテーテル類の安全を確認"],
            tc=PC_NAVY,bg=PC_LIGHT)
    infobox(sl,Inches(6.45),Inches(2.15),Inches(5.9),Inches(4.95),
            "サポート役（足側）",
            ["✓  患者の下肢・臀部を支持",
             "✓  リーダーの指示に従う",
             "✓  足先から骨盤まで確実に保持",
             "✓  移動方向を確認してから動く"],
            tc=PC_TEAL,bg=PC_GRN_BG,bc=PC_TEAL)
    imgbox(sl,Inches(0.35),Inches(7.18),Inches(12.63),Inches(0.25),
           "【画像⑨】2名介助でのベッド上移動・移乗イラスト（頭側・足側の持ち方）")

    # ── Slide 13: 体位変換 ───────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"体位変換（臥位での移動）","褥瘡予防と安楽な体位管理","13 / 18")
    txbox(sl,Inches(0.35),Inches(1.45),Inches(12.63),Inches(0.4),
          "体位変換は原則2時間ごと。骨突出部（仙骨・踵・大転子等）への圧迫を避けること",
          11,color=PC_GRAY)
    positions13=[
        (PC_NAVY, "仰臥位","（あお向け）","最も基本の体位\n踵・仙骨の圧迫に注意","【画像⑩-1】仰臥位"),
        (PC_TEAL, "側臥位","（横向き）","30°傾けが理想\nクッションで保持","【画像⑩-2】側臥位"),
        (PC_AMBER,"半座位","（ファウラー位）","床頭台15〜45°\n誤嚥予防に有効","【画像⑩-3】ファウラー位"),
        (PC_GREEN,"腹臥位","（うつぶせ）","呼吸改善に有効\n顔・腹部の圧迫注意","【画像⑩-4】腹臥位"),
    ]
    for i,(c,t,sub,d,img) in enumerate(positions13):
        x=Inches(0.35+i*3.2)
        imgbox(sl,x,Inches(1.9),Inches(3.0),Inches(2.2),img)
        rect(sl,x,Inches(4.15),Inches(3.0),Inches(2.6),fill=PC_BG,line=c,lw=19050)
        txbox(sl,x+Inches(0.1),Inches(4.2),Inches(2.8),Inches(0.5),f"{t}  {sub}",13,True,c)
        txbox(sl,x+Inches(0.1),Inches(4.78),Inches(2.8),Inches(1.8),d,11,color=PC_DARK)
    warnbox(sl,Inches(0.35),Inches(6.85),Inches(12.63),Inches(0.58),
            "体位変換のたびに皮膚状態・褥瘡チェックを必ず実施！ポジショニングクッションで圧迫を分散")

    # ── Slide 14: 転倒・転落予防 ─────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"転倒・転落予防","リスク評価と予防策の徹底","14 / 18")
    infobox(sl,Inches(0.35),Inches(1.55),Inches(3.8),Inches(5.0),
            "患者側リスク",
            ["▸ 筋力低下・バランス障害",
             "▸ 認知症・見当識障害",
             "▸ 薬剤影響（睡眠薬・降圧薬等）",
             "▸ 起立性低血圧",
             "▸ 排泄に関する切迫感"],
            tc=PC_RED,bg=PC_RED_BG,bc=PC_RED)
    infobox(sl,Inches(4.35),Inches(1.55),Inches(3.8),Inches(5.0),
            "環境側リスク",
            ["▸ 床の濡れ・障害物",
             "▸ 不適切な照明（暗すぎ）",
             "▸ ベッド高さ不適切",
             "▸ ナースコール未設置",
             "▸ スリッパなど滑りやすい履物"],
            tc=PC_AMBER,bg=PC_WARN_BG,bc=PC_AMBER)
    infobox(sl,Inches(8.35),Inches(1.55),Inches(4.63),Inches(5.0),
            "予防策",
            ["▸ 転倒リスクスコア評価（毎日）",
             "▸ 離床センサー・ベッドアラーム設置",
             "▸ 低床ベッドの使用",
             "▸ ナースコール手の届く位置へ",
             "▸ ラウンド強化（排泄パターン把握）"],
            tc=PC_GREEN,bg=PC_GRN_BG,bc=PC_GREEN)
    warnbox(sl,Inches(0.35),Inches(6.65),Inches(12.63),Inches(0.77),
            "転倒・転落発生時：患者の安全確保 → ナースコール → 状態確認 → 報告 → 記録",danger=True)

    # ── Slide 15: よくあるミス ───────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"よくあるミスとその対策","新人職員が陥りやすいポイント","15 / 18")
    # テーブル
    cols_w=[Inches(3.0),Inches(4.2),Inches(5.4)]
    hdrs=["よくあるミス","起こりうる問題","対策"]
    hdr_colors=[PC_RED,PC_AMBER,PC_GREEN]
    rows15=[
        ("✗ 声かけを忘れる","患者の心理的不安・筋緊張が増す","✓ 必ず「〇〇しますね」と伝えてから開始"),
        ("✗ ブレーキを忘れる","移乗中に車椅子が動いて転落","✓ 移乗前の確認を手順書通りに実施"),
        ("✗ 腰だけで持ち上げる","介助者の腰痛・ぎっくり腰","✓ ボディメカニクスを意識。足・体幹を使う"),
        ("✗ 急ぎすぎる","患者が不安・バランス崩す","✓ 患者のペースに合わせ、ゆっくり確実に"),
        ("✗ 1人で無理をする","患者・介助者ともに受傷リスク","✓ 迷ったら必ず同僚に声をかける"),
        ("✗ カテーテルを忘れる","ドレーン抜去・皮膚損傷","✓ 移乗前にルート類の長さ・固定を確認"),
    ]
    for j,(h,c) in enumerate(zip(hdrs,hdr_colors)):
        x=sum(cols_w[:j]) if j>0 else Inches(0.35)
        if j>0: x = Inches(0.35)+sum(cols_w[:j])
        rect(sl,x,Inches(1.55),cols_w[j],Inches(0.55),fill=c)
        txbox(sl,x+Inches(0.08),Inches(1.6),cols_w[j]-Inches(0.1),Inches(0.45),h,13,True,PC_WHITE)
    for i,row in enumerate(rows15):
        y=Inches(2.1)+i*Inches(0.83)
        bg=(0xF9,0xFA,0xFB) if i%2==0 else (0xFF,0xFF,0xFF)
        x=Inches(0.35)
        for j,(val,c) in enumerate(zip(row,hdr_colors)):
            w_=cols_w[j]
            rect(sl,x,y,w_,Inches(0.78),fill=bg,line=(0xD1,0xD5,0xDB),lw=9525)
            col_=(0xB9,0x1C,0x1C) if j==0 else (0x09,0x17,0x2D) if j==2 else PC_DARK
            txbox(sl,x+Inches(0.08),y+Inches(0.04),w_-Inches(0.12),Inches(0.7),val,10,color=col_)
            x+=w_
    warnbox(sl,Inches(0.35),Inches(7.08),Inches(12.63),Inches(0.35),
            "わからないことは必ず先輩に確認！報告・連絡・相談を徹底しましょう")

    # ── Slide 16: 記録・報告 ─────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"記録・報告のポイント","介助後の観察と情報共有","16 / 18")
    items16=[
        (PC_NAVY, "いつ","移乗・移動の実施日時"),
        (PC_BLUE, "どこで","場所（病室・廊下・リハビリ室など）"),
        (PC_TEAL, "何を","実施した介助内容（車椅子移乗・歩行介助など）"),
        (PC_GREEN,"誰が","実施者（1名・2名介助など）"),
        (PC_AMBER,"どのように","使用した補助具・患者の協力度"),
        (PC_PURPLE,"患者の状態","バイタル変化・表情・疼痛・ふらつきの有無"),
    ]
    for i,(c,t,d) in enumerate(items16):
        col=i%3; row=i//3
        x=Inches(0.35+col*4.1); y=Inches(1.55)+row*Inches(1.55)
        rect(sl,x,y,Inches(3.9),Inches(1.4),fill=PC_BG,line=c,lw=19050)
        s=oval(sl,x+Inches(0.12),y+Inches(0.12),Inches(0.5),Inches(0.5),fill=c)
        p=s.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=t; r.font.size=Pt(9); r.font.bold=True; r.font.color.rgb=PC_WHITE
        txbox(sl,x+Inches(0.72),y+Inches(0.12),Inches(3.06),Inches(1.15),d,11,color=PC_DARK)
    # 報告必要な場合
    rect(sl,Inches(0.35),Inches(4.75),Inches(12.63),Inches(2.55),fill=PC_RED_BG,line=PC_RED,lw=25400)
    txbox(sl,Inches(0.5),Inches(4.82),Inches(12.0),Inches(0.45),
          "報告が必要な場合",14,True,PC_RED)
    report_items=["転倒・転落発生","ヒヤリハット","患者が痛みを訴えた",
                  "バイタルが基準値を逸脱","皮膚の発赤・損傷発見","ドレーン・カテーテルが抜けた","患者が介助を拒否"]
    for i,it in enumerate(report_items):
        col=i%4; row_=i//4
        x=Inches(0.5+col*3.1); y=Inches(5.35)+row_*Inches(0.7)
        txbox(sl,x,y,Inches(2.95),Inches(0.55),f"⚠  {it}",11,color=PC_DARK)
    warnbox(sl,Inches(0.35),Inches(7.1),Inches(12.63),Inches(0.33),
            "記録は「事実」を具体的・客観的に。主観的な表現は避ける")

    # ── Slide 17: 確認テスト ─────────────────────────
    sl=prs.slides.add_slide(blank)
    header(sl,"確認テスト","本日の研修内容を振り返りましょう","17 / 18",bg=PC_TEAL)
    questions=[
        ("Q1","ボディメカニクスの目的は何ですか？\n①患者の安楽  ②介助者の腰痛予防  ③両方  ④どちらでもない"),
        ("Q2","片麻痺患者への車椅子移乗で、車椅子を置く位置はどちら側ですか？\n①患側（麻痺側）  ②健側（麻痺のない側）  ③どちらでも同じ"),
        ("Q3","2名介助が必要な条件を1つ答えてください。"),
        ("Q4","移乗介助前に必ず確認すべき3つのことを答えてください。"),
        ("Q5","転倒が発生した際の対応の優先順位を答えてください。"),
    ]
    for i,(q,text) in enumerate(questions):
        col=i%3; row_=i//3
        x=Inches(0.35+col*4.3); y=Inches(1.55)+row_*Inches(2.5)
        h_=Inches(2.35) if row_==0 else Inches(2.55)
        rect(sl,x,y,Inches(4.1),h_,fill=PC_BG,line=PC_TEAL,lw=25400)
        s=oval(sl,x+Inches(1.55),y+Inches(0.1),Inches(0.6),Inches(0.6),fill=PC_TEAL)
        p=s.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=q; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=PC_WHITE
        txbox(sl,x+Inches(0.12),y+Inches(0.82),Inches(3.86),h_-Inches(0.95),text,11,color=PC_DARK)
    warnbox(sl,Inches(0.35),Inches(7.1),Inches(12.63),Inches(0.35),
            "答え合わせは研修担当者と一緒に行いましょう！")

    # ── Slide 18: まとめ ─────────────────────────────
    sl=prs.slides.add_slide(blank)
    set_bg(sl,PC_NAVY)
    rect(sl,0,0,W,H,fill=PC_NAVY)
    rect(sl,Inches(9.5),0,Inches(3.83),H,fill=(0x1D,0x4E,0xD8))
    txbox(sl,Inches(0.5),Inches(0.25),Inches(9.0),Inches(0.75),
          "本日のまとめ",28,True,PC_WHITE)
    txbox(sl,Inches(0.5),Inches(0.95),Inches(9.0),Inches(0.4),
          "移乗・移動介助のポイントを復習",14,color=PC_LNAV)
    summary=[
        (PC_LBLUE,  "1","介助の3原則","安全・安楽・自立支援を常に意識する"),
        (PC_TEAL,   "2","ボディメカニクス","正しい姿勢で患者も介助者も守る"),
        (PC_GREEN,  "3","確認・声かけ","介助前の確認と患者への声かけを徹底する"),
        (PC_AMBER,  "4","移乗手順の遵守","車椅子の角度・ブレーキ・フットレストを確認"),
        (PC_RED,    "5","2名介助の判断","一人で無理せず、迷ったら必ず声をかける"),
        (PC_PURPLE, "6","記録・報告","気づいたことはすぐに記録・報告・連絡"),
    ]
    for i,(c,n,t,d) in enumerate(summary):
        col=i%3; row_=i//3
        x=Inches(0.5+col*3.0); y=Inches(1.5)+row_*Inches(2.3)
        rect(sl,x,y,Inches(2.85),Inches(2.1),fill=(0x1E,0x40,0xAF),line=c,lw=25400)
        s=oval(sl,x+Inches(1.1),y+Inches(0.1),Inches(0.62),Inches(0.62),fill=c)
        p=s.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=n; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=PC_WHITE
        txbox(sl,x+Inches(0.1),y+Inches(0.85),Inches(2.65),Inches(0.5),t,13,True,c,align=PP_ALIGN.CENTER)
        txbox(sl,x+Inches(0.1),y+Inches(1.38),Inches(2.65),Inches(0.6),d,10,color=PC_LNAV)
    txbox(sl,Inches(0.5),Inches(7.1),Inches(9.0),Inches(0.35),
          "患者の「安全・安楽・自立支援」のために、今日から実践しましょう！",
          12,True,PC_LNAV)

    prs.save("C:/Users/yuuna/agens/新人勉強会/スライド_v2_移乗移動介助.pptx")
    print("OK: PPTX saved")


# ════════════════════════════════════════════════════════
#  Word 配布資料（白黒）
# ════════════════════════════════════════════════════════
def set_cell_bg(cell,r,g,b):
    tc=cell._tc; tcPr=tc.get_or_add_tcPr()
    shd=OxmlElement('w:shd')
    shd.set(qn('w:val'),'clear'); shd.set(qn('w:color'),'auto')
    shd.set(qn('w:fill'),f'{r:02X}{g:02X}{b:02X}')
    tcPr.append(shd)

def set_para_bg(para,r,g,b):
    pPr=para._p.get_or_add_pPr()
    shd=OxmlElement('w:shd')
    shd.set(qn('w:val'),'clear'); shd.set(qn('w:color'),'auto')
    shd.set(qn('w:fill'),f'{r:02X}{g:02X}{b:02X}')
    pPr.append(shd)

BK=DRGBColor(0x00,0x00,0x00)
DG=DRGBColor(0x37,0x41,0x51)
MG=DRGBColor(0x6B,0x72,0x80)
LG=DRGBColor(0xE5,0xE7,0xEB)
VL=DRGBColor(0xF9,0xFA,0xFB)
WH=DRGBColor(0xFF,0xFF,0xFF)

def d_heading(doc,text,level=1,size=16):
    p=doc.add_paragraph()
    if level==1: set_para_bg(p,0x37,0x41,0x51)
    elif level==2: set_para_bg(p,0xE5,0xE7,0xEB)
    r=p.add_run(text); r.bold=True; r.font.size=DPt(size)
    r.font.color.rgb = WH if level==1 else BK
    p.paragraph_format.space_before=DPt(6)
    p.paragraph_format.space_after=DPt(4)
    return p

def d_body(doc,text,size=9.5):
    p=doc.add_paragraph()
    r=p.add_run(text); r.font.size=DPt(size); r.font.color.rgb=DG
    p.paragraph_format.space_after=DPt(3)
    return p

def d_bullet(doc,text,size=9.5):
    p=doc.add_paragraph(style='List Bullet')
    p.paragraph_format.left_indent=DCm(0.4)
    r=p.add_run(text); r.font.size=DPt(size); r.font.color.rgb=DG
    p.paragraph_format.space_after=DPt(2)

def d_step_table(doc,steps):
    tbl=doc.add_table(rows=1,cols=3); tbl.style='Table Grid'
    tbl.alignment=WD_TABLE_ALIGNMENT.CENTER
    hdrs=["No.","手順・内容","ポイント・注意"]
    for i,(h_,c) in enumerate(zip(hdrs,tbl.rows[0].cells)):
        c.text=h_; c.paragraphs[0].runs[0].bold=True
        c.paragraphs[0].runs[0].font.size=DPt(9)
        c.paragraphs[0].runs[0].font.color.rgb=WH
        set_cell_bg(c,0x37,0x41,0x51)
    for i,c_w in enumerate([DCm(1.0),DCm(7.5),DCm(7.0)]):
        for cell in tbl.columns[i].cells: cell.width=c_w
    for i,row in enumerate(steps):
        cells=tbl.add_row().cells
        cells[0].text=str(row[0])
        cells[0].paragraphs[0].runs[0].font.size=DPt(10)
        cells[0].paragraphs[0].runs[0].bold=True
        cells[1].text=row[1]; cells[1].paragraphs[0].runs[0].font.size=DPt(9)
        cells[2].text=row[2]; cells[2].paragraphs[0].runs[0].font.size=DPt(9)
        cells[2].paragraphs[0].runs[0].font.color.rgb=MG
        if i%2==0:
            for c in cells: set_cell_bg(c,0xF9,0xFA,0xFB)
    return tbl

def d_warn_box(doc,text,icon="▲"):
    tbl=doc.add_table(rows=1,cols=1); tbl.style='Table Grid'
    cell=tbl.cell(0,0)
    cell.text=f"{icon}  {text}"
    cell.paragraphs[0].runs[0].font.size=DPt(9)
    cell.paragraphs[0].runs[0].bold=True
    set_cell_bg(cell,0xE5,0xE7,0xEB)
    return tbl

def d_check_list(doc,items,label="チェックリスト"):
    p=doc.add_paragraph()
    r=p.add_run(f"□ {label}"); r.bold=True; r.font.size=DPt(10)
    for it in items:
        p2=doc.add_paragraph(style='List Bullet')
        p2.paragraph_format.left_indent=DCm(0.5)
        r2=p2.add_run(f"□ {it}"); r2.font.size=DPt(9)

def d_2col_table(doc,rows_data,h1,h2):
    tbl=doc.add_table(rows=1,cols=2); tbl.style='Table Grid'
    for cell,h in zip(tbl.rows[0].cells,[h1,h2]):
        cell.text=h; cell.paragraphs[0].runs[0].bold=True
        cell.paragraphs[0].runs[0].font.size=DPt(9)
        cell.paragraphs[0].runs[0].font.color.rgb=WH
        set_cell_bg(cell,0x37,0x41,0x51)
    for cell in tbl.columns[0].cells: cell.width=DCm(7.5)
    for cell in tbl.columns[1].cells: cell.width=DCm(8.0)
    for i,(l,r) in enumerate(rows_data):
        cells=tbl.add_row().cells
        cells[0].text=l; cells[0].paragraphs[0].runs[0].font.size=DPt(9)
        cells[1].text=r; cells[1].paragraphs[0].runs[0].font.size=DPt(9)
        if i%2==0:
            for c in cells: set_cell_bg(c,0xF9,0xFA,0xFB)
    return tbl

def make_word():
    doc=Document()
    for sec in doc.sections:
        sec.top_margin=DCm(2.0); sec.bottom_margin=DCm(2.0)
        sec.left_margin=DCm(2.2); sec.right_margin=DCm(2.2)

    # 表紙
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
    set_para_bg(p,0x37,0x41,0x51)
    r=p.add_run("新人職員研修　配布資料")
    r.bold=True; r.font.size=DPt(10); r.font.color.rgb=DRGBColor(0xD1,0xD5,0xDB)

    p2=doc.add_paragraph(); p2.alignment=WD_ALIGN_PARAGRAPH.CENTER
    set_para_bg(p2,0x37,0x41,0x51)
    r2=p2.add_run("移乗・移動介助の基本")
    r2.bold=True; r2.font.size=DPt(24); r2.font.color.rgb=WH

    p3=doc.add_paragraph(); p3.alignment=WD_ALIGN_PARAGRAPH.CENTER
    set_para_bg(p3,0x37,0x41,0x51)
    r3=p3.add_run("〜安全・安楽・自立支援を実践する〜")
    r3.font.size=DPt(11); r3.font.color.rgb=DRGBColor(0xD1,0xD5,0xDB)

    p4=doc.add_paragraph(); p4.alignment=WD_ALIGN_PARAGRAPH.CENTER
    set_para_bg(p4,0x37,0x41,0x51)
    p4.paragraph_format.space_after=DPt(8)
    r4=p4.add_run("対象：看護職員・介護職員　｜　所要時間：約30分　｜　医療法人 研修・教育委員会")
    r4.font.size=DPt(9); r4.font.color.rgb=DRGBColor(0x9C,0xA3,0xAF)

    # 受講者情報
    info_tbl=doc.add_table(rows=1,cols=3); info_tbl.style='Table Grid'
    for cell,txt in zip(info_tbl.rows[0].cells,
                        [" 実施日：                 "," 担当講師：                 "," 受講者名：                 "]):
        cell.text=txt; cell.paragraphs[0].runs[0].font.size=DPt(9)
        set_cell_bg(cell,0xF3,0xF4,0xF6)
    doc.add_paragraph()

    # ── 研修目標 ────────────────────────────────
    d_heading(doc,"研修目標",1,13)
    for n,t in [("①","移乗・移動介助の基本原則とボディメカニクスを理解する"),
                ("②","車椅子・ベッド間の移乗手順を正しく実施できる"),
                ("③","歩行介助・立ち上がり介助を安全に行える"),
                ("④","転倒・転落リスクを予測し予防策を講じられる"),
                ("⑤","患者の尊厳と残存能力を尊重した介助ができる")]:
        p=doc.add_paragraph(style='List Bullet')
        p.paragraph_format.left_indent=DCm(0.5)
        r=p.add_run(f"{n}  {t}"); r.font.size=DPt(9.5)
    doc.add_paragraph()

    # ── 介助の3原則 ────────────────────────────
    d_heading(doc,"1. 介助の3原則",1,13)
    tbl3=doc.add_table(rows=1,cols=3); tbl3.style='Table Grid'
    for cell,h in zip(tbl3.rows[0].cells,["安　全","安　楽","自立支援"]):
        cell.text=h; cell.paragraphs[0].runs[0].bold=True
        cell.paragraphs[0].runs[0].font.size=DPt(11)
        cell.paragraphs[0].runs[0].font.color.rgb=WH
        set_cell_bg(cell,0x37,0x41,0x51)
        cell.paragraphs[0].alignment=WD_ALIGN_PARAGRAPH.CENTER
    rows3=[
        ("転倒・転落・皮膚損傷を防ぐ","患者の苦痛・不安を最小化","できることは自分で行ってもらう"),
        ("環境整備を必ず行う","声かけ・同意を必ず得る","残存能力を活かす介助"),
        ("2人介助の基準を遵守する","プライバシーを守る","ADL維持・向上を目指す"),
    ]
    for i,row in enumerate(rows3):
        cells=tbl3.add_row().cells
        for j,val in enumerate(row):
            cells[j].text=f"✓  {val}"; cells[j].paragraphs[0].runs[0].font.size=DPt(9)
            if i%2==0: set_cell_bg(cells[j],0xF9,0xFA,0xFB)
    doc.add_paragraph()

    # ── ボディメカニクス ──────────────────────────
    d_heading(doc,"2. ボディメカニクスの基本",1,13)
    d_body(doc,"人体の構造と物理の法則を活かし、介助者の体への負担を最小限にする技術です。腰痛は職業病ではありません。")
    bm_steps=[
        (1,"足を肩幅に開く（支持基底面を広げる）","足の間隔が広いほど安定。安定した姿勢を確保する"),
        (2,"膝を曲げて重心を低くする","膝を曲げて腰を落とし、重心を下げる"),
        (3,"重心を近づける","患者との距離を縮め、密着して介助する"),
        (4,"大きな筋群を使う","腰だけでなく脚・体幹の大きな筋肉を使う"),
        (5,"身体をひとつにまとめる","患者の手足を体に引き付けてから動かす"),
        (6,"ねじり動作を避ける","足を向けてから体を回旋させる。ひねり禁止"),
    ]
    d_step_table(doc,bm_steps)
    doc.add_paragraph()

    # ── 介助前確認 ───────────────────────────────
    d_heading(doc,"3. 介助前の確認事項",1,13)
    d_heading(doc,"患者確認",2,10)
    for it in ["患者の理解・同意を得る","バイタルサイン・体調確認",
               "点滴・ドレーン・カテーテル確認","麻痺・疼痛・可動域の確認","患者の体重・介助レベル確認"]:
        d_bullet(doc,it)
    d_heading(doc,"環境整備",2,10)
    for it in ["ベッドの高さ調整（介助しやすい高さ）","ブレーキ・ストッパーの確認",
               "床の濡れ・障害物の除去","スリッパ→滑り止め靴に変更","プライバシーカーテンを閉める"]:
        d_bullet(doc,it)
    d_warn_box(doc,"2人介助基準：全介助・体重60kg以上・不安定な患者は必ず2名で対応","!")
    doc.add_paragraph()
    doc.add_page_break()

    # ── 立ち上がり介助 ───────────────────────────
    d_heading(doc,"4. 立ち上がり介助（端座位 → 立位）",1,13)
    d_step_table(doc,[
        (1,"端座位を確保する","ベッド端に座り、足底を床につける"),
        (2,"前傾姿勢をとる（鼻がひざの上）","体を前に傾け、重心を前方に移動"),
        (3,"「せーの」の声かけで一緒に立ち上がる","患者の脇・腰を支えながら立位へ"),
        (4,"立位の安定確認","めまい・ふらつきがないか確認。必要なら手すりを"),
    ])
    doc.add_paragraph()

    # ── ベッド→車椅子 ───────────────────────────
    d_heading(doc,"5. 移乗介助① ベッド → 車椅子",1,13)
    d_body(doc,"片麻痺患者の場合：健側に車椅子を置く。「ブレーキ確認」を声に出して習慣化する。")
    d_step_table(doc,[
        (1,"車椅子をベッドに対して約30°の角度で設置（健側に置く）","ブレーキON・フットレストを上げる"),
        (2,"患者を端座位にし、足底を床につける","介助者は患側に立つ"),
        (3,"前傾姿勢をとってもらう","膝を介助者の膝でブロック"),
        (4,"「せーの」で一緒に立ち上がる","腰〜臀部を支える。脇を持つのはNG"),
        (5,"ゆっくり方向転換し、車椅子の方向へ","足をすり足で動かす。ねじり禁止"),
        (6,"臀部をシートに誘導しながら着座させる","深く座れているか確認"),
        (7,"フットレストを下げ、足を乗せる","姿勢を整えて「座れましたか？」と声かけ"),
        (8,"安全確認・記録","ブレーキ解除・移動可能か確認"),
    ])
    doc.add_paragraph()

    # ── 車椅子→ベッド ───────────────────────────
    d_heading(doc,"6. 移乗介助② 車椅子 → ベッド",1,13)
    d_body(doc,"基本は逆手順。健側から30°配置 → ブレーキ・フットレスト確認 → 前傾 → 立位 → 方向転換 → 着座。")
    d_2col_table(doc,[
        ("車椅子の配置","ベッドに対して健側に30°。ブレーキ・フットレスト確認を先に行う"),
        ("体重移動の誘導","「前に体を倒して」と声かけ。前傾姿勢から立ち上がり誘導"),
        ("方向転換","介助者は患者の患側に位置。体を密着させる"),
        ("着座誘導","ベッドに手をつかせてから、ゆっくり膝を曲げて座らせる"),
        ("体位の整え","座位バランス確認後、臥位等へ体位を整える"),
    ],"ポイント","詳細")
    doc.add_paragraph()

    # ── 移乗補助用具 ─────────────────────────────
    d_heading(doc,"7. 移乗補助用具の活用",1,13)
    d_2col_table(doc,[
        ("スライディングボード","座位移乗（車椅子⇔ベッド・シャワーチェア）。麻痺・筋力低下患者に有効。使用時は2名確認"),
        ("移乗ベルト（トランスファーベルト）","患者の腰に装着。急な体重移動に対応。転倒防止に有効"),
        ("介護リフト","全介助・重度障害患者に使用。スリングシートを正しく装着。操作は研修後のみ"),
    ],"補助具","特徴・使用方法")
    doc.add_paragraph()
    doc.add_page_break()

    # ── 歩行介助 ─────────────────────────────────
    d_heading(doc,"8. 歩行介助の基本",1,13)
    d_body(doc,"介助者の立ち位置：患側（麻痺のある側）の斜め後方。片手は患者の腰または移乗ベルトを握る。")
    d_step_table(doc,[
        (1,"開始前","靴・履物確認。バイタル確認。歩行補助具の用意"),
        (2,"出発時","ベッドサイドで立位保持。めまいがないか確認"),
        (3,"歩行中","声かけしながら一定ペースで。障害物を先に伝える"),
        (4,"方向転換","健側に向かってゆっくり回転。急な向き変え禁止"),
        (5,"着座時","椅子/ベッドに確実に触れてからゆっくり着座"),
    ])
    d_warn_box(doc,"脇を抱えて支えるのはNG！肩関節脱臼のリスクがあります。腰・骨盤を支えましょう。","!")
    doc.add_paragraph()

    # ── 歩行補助具 ───────────────────────────────
    d_heading(doc,"9. 歩行補助具の種類と特徴",1,13)
    d_2col_table(doc,[
        ("T字杖（ステッキ）","軽度の歩行障害に適応。健側で使用。4点杖より安定性は低い"),
        ("4点杖","接地面が広く安定。片麻痺・バランス障害に有効。室内使用が主"),
        ("歩行器（固定型）","前方支持で安定。両上肢に体重をかけられる。段差・屋外には不向き"),
        ("ロールウォーカー","車輪付きで歩きやすい。屋内外兼用可。ブレーキ操作を必ず確認"),
    ],"補助具","特徴")
    d_body(doc,"補助具の選択は理学療法士・主治医と連携して決定すること。")
    doc.add_paragraph()

    # ── 全介助（2名） ────────────────────────────
    d_heading(doc,"10. 全介助の移乗（2名介助）",1,13)
    d_warn_box(doc,"2名介助の適応：全介助患者 ／ 体重60kg以上 ／ 不安定な患者 ／ 初回移乗","!")
    doc.add_paragraph()
    d_2col_table(doc,[
        ("リーダー役（頭側）",
         "上半身・頭頸部を支持。声かけ・タイミングの指示を出す。「いちにのさん」で全員同時に動く。気道・カテーテル類の安全を確認"),
        ("サポート役（足側）",
         "下肢・臀部を支持。リーダーの指示に従う。足先から骨盤まで確実に保持。移動方向を確認してから動く"),
    ],"役割","内容")
    doc.add_paragraph()

    # ── 体位変換 ─────────────────────────────────
    d_heading(doc,"11. 体位変換（臥位での移動）",1,13)
    d_body(doc,"体位変換は原則2時間ごと。骨突出部（仙骨・踵・大転子等）への圧迫を避ける。")
    d_2col_table(doc,[
        ("仰臥位（あお向け）","最も基本の体位。踵・仙骨の圧迫に注意"),
        ("側臥位（横向き）","30°傾けが理想。クッションで保持"),
        ("半座位（ファウラー位）","床頭台15〜45°。誤嚥予防に有効"),
        ("腹臥位（うつぶせ）","呼吸改善に有効。顔・腹部の圧迫注意"),
    ],"体位","説明・注意点")
    d_warn_box(doc,"体位変換のたびに皮膚状態・褥瘡チェックを必ず実施！ポジショニングクッションを活用する。")
    doc.add_paragraph()
    doc.add_page_break()

    # ── 転倒・転落予防 ───────────────────────────
    d_heading(doc,"12. 転倒・転落予防",1,13)
    d_2col_table(doc,[
        ("患者側リスク","筋力低下・バランス障害、認知症、薬剤影響（睡眠薬・降圧薬等）、起立性低血圧、排泄の切迫感"),
        ("環境側リスク","床の濡れ・障害物、不適切な照明、ベッド高さ不適切、ナースコール未設置、滑りやすい履物"),
        ("予防策","転倒リスクスコア評価（毎日）、離床センサー設置、低床ベッド使用、ナースコール手の届く位置、ラウンド強化"),
    ],"リスク・予防の種類","内容")
    d_warn_box(doc,"転倒・転落発生時：患者の安全確保 → ナースコール → 状態確認 → 報告 → 記録","!")
    doc.add_paragraph()

    # ── よくあるミスと対策 ───────────────────────
    d_heading(doc,"13. よくあるミスとその対策",1,13)
    d_step_table(doc,[
        ("✗","声かけを忘れる → 患者の不安・筋緊張増","必ず「〇〇しますね」と伝えてから開始"),
        ("✗","ブレーキを忘れる → 転落","移乗前の確認を手順書通りに実施"),
        ("✗","腰だけで持ち上げる → 介助者腰痛","ボディメカニクスを意識。足・体幹を使う"),
        ("✗","急ぎすぎる → バランス崩す","患者のペースに合わせ、ゆっくり確実に"),
        ("✗","1人で無理をする → 双方受傷","迷ったら必ず同僚に声をかける"),
        ("✗","カテーテルを忘れる → 抜去","移乗前にルート類の長さ・固定を確認"),
    ])
    doc.add_paragraph()

    # ── 記録・報告 ───────────────────────────────
    d_heading(doc,"14. 記録・報告のポイント",1,13)
    d_body(doc,"5W1H＋患者の状態を記録する。")
    d_2col_table(doc,[
        ("いつ","移乗・移動の実施日時"),
        ("どこで","場所（病室・廊下・リハビリ室など）"),
        ("何を","実施した介助内容（車椅子移乗・歩行介助など）"),
        ("誰が","実施者（1名・2名介助など）"),
        ("どのように","使用した補助具・患者の協力度"),
        ("患者の状態","バイタル変化・表情・疼痛・ふらつきの有無"),
    ],"記録項目","内容")
    doc.add_paragraph()
    d_heading(doc,"報告が必要な場合",2,10)
    for it in ["転倒・転落が発生した","ヒヤリハットがあった","患者が痛みを訴えた",
               "バイタルが基準値を逸脱した","皮膚の発赤・損傷を発見した",
               "ドレーン・カテーテルが抜けた","患者が介助を拒否した"]:
        d_bullet(doc,f"▲  {it}")
    doc.add_paragraph()

    # ── 確認テスト ───────────────────────────────
    d_heading(doc,"15. 確認テスト",1,13)
    d_body(doc,"研修後に自分で答えを記入してみましょう。担当者と一緒に答え合わせを行います。")
    qs=[
        ("Q1","ボディメカニクスの目的は何ですか？\n①患者の安楽  ②介助者の腰痛予防  ③両方  ④どちらでもない"),
        ("Q2","片麻痺患者への車椅子移乗で、車椅子を置く位置はどちら側ですか？\n①患側（麻痺側）  ②健側（麻痺のない側）  ③どちらでも同じ"),
        ("Q3","2名介助が必要な条件を1つ答えてください。"),
        ("Q4","移乗介助前に必ず確認すべき3つのことを答えてください。"),
        ("Q5","転倒が発生した際の対応の優先順位を答えてください。"),
    ]
    for q,text in qs:
        p=doc.add_paragraph()
        r=p.add_run(q); r.bold=True; r.font.size=DPt(10)
        p.paragraph_format.space_before=DPt(4)
        p2=doc.add_paragraph()
        r2=p2.add_run(text); r2.font.size=DPt(9.5)
        p3=doc.add_paragraph()
        r3=p3.add_run("回答：                                                                                    ")
        r3.font.size=DPt(9.5); r3.font.color.rgb=MG
    doc.add_page_break()

    # ── まとめ・チェックリスト ──────────────────
    d_heading(doc,"16. まとめ & 自己チェックリスト",1,13)
    tbl_s=doc.add_table(rows=1,cols=2); tbl_s.style='Table Grid'
    for cell,h in zip(tbl_s.rows[0].cells,["#  重要ポイント","実践内容"]):
        cell.text=h; cell.paragraphs[0].runs[0].bold=True
        cell.paragraphs[0].runs[0].font.size=DPt(9)
        cell.paragraphs[0].runs[0].font.color.rgb=WH
        set_cell_bg(cell,0x37,0x41,0x51)
    for cell in tbl_s.columns[0].cells: cell.width=DCm(3.5)
    for cell in tbl_s.columns[1].cells: cell.width=DCm(12.0)
    summaries=[
        ("1. 介助の3原則","安全・安楽・自立支援を常に意識する"),
        ("2. ボディメカニクス","正しい姿勢で患者も介助者も守る"),
        ("3. 確認・声かけ","介助前の確認と患者への声かけを徹底する"),
        ("4. 移乗手順の遵守","車椅子の角度・ブレーキ・フットレストを確認"),
        ("5. 2名介助の判断","一人で無理せず、迷ったら必ず声をかける"),
        ("6. 記録・報告","気づいたことはすぐに記録・報告・連絡"),
    ]
    for i,(k,v) in enumerate(summaries):
        cells=tbl_s.add_row().cells
        cells[0].text=k; cells[0].paragraphs[0].runs[0].font.size=DPt(9)
        cells[0].paragraphs[0].runs[0].bold=True
        cells[1].text=v; cells[1].paragraphs[0].runs[0].font.size=DPt(9)
        if i%2==0:
            for c in cells: set_cell_bg(c,0xF9,0xFA,0xFB)
    doc.add_paragraph()
    d_check_list(doc,[
        "ボディメカニクス6原則をすらすら言えますか？",
        "ベッド→車椅子の移乗手順を8ステップで説明できますか？",
        "ブレーキ確認を声に出す習慣がつきましたか？",
        "患側（麻痺側）の斜め後ろに立って歩行介助できますか？",
        "段差での車椅子後ろ向き操作を理解しましたか？",
        "転倒リスク5つを挙げられますか？",
        "ヒヤリハットを必ず報告するという意識が持てましたか？",
    ],"研修後 自己チェックリスト")
    doc.add_paragraph()

    # フッター
    p_f=doc.add_paragraph(); p_f.alignment=WD_ALIGN_PARAGRAPH.CENTER
    set_para_bg(p_f,0x37,0x41,0x51)
    r_f=p_f.add_run("新人職員研修｜移乗・移動介助の基本　配布資料　©研修・教育委員会")
    r_f.font.size=DPt(8); r_f.font.color.rgb=DRGBColor(0x9C,0xA3,0xAF)

    doc.save("C:/Users/yuuna/agens/新人勉強会/配布資料_白黒_v2.docx")
    print("OK: Word B&W saved")


if __name__=="__main__":
    make_pptx()
    make_word()
    print("DONE")
